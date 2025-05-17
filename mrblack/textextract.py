#!/usr/bin/env python3
# -*- coding: utf-8 -*-
#
# File: textextract.py
# Author: Wadih Khairallah
# Created: 2024-12-01 12:12:08
# Modified: 2025-05-17 19:04:56
# Enhanced with additional features

import os
import re
import sys
import pwd
import json
import math
import socket
import platform
import subprocess
import hashlib
import string
import tempfile
import random
import unicodedata
import concurrent.futures
import logging
import shutil
from logging import Logger
from pathlib import Path

from uuid import uuid4
from datetime import datetime
from io import StringIO
from typing import (
    Optional,
    Dict,
    Any,
    List,
    Union,
    Tuple,
    Set,
    Callable,
)

# Web Specific
import requests
from requests_html import HTMLSession
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup

# Text Specific
from readability import Document as RDocument
import magic
import pytesseract
import pandas as pd
import speech_recognition as sr
import pymupdf
from docx import Document
from mss import mss
from pydub import AudioSegment
from deep_translator import GoogleTranslator

# Image Specific
from PIL import Image

# NLP and Text Analytics
from collections import Counter
import nltk
from nltk.corpus import stopwords, wordnet
from nltk.stem import WordNetLemmatizer, PorterStemmer
from nltk.util import ngrams
from nltk.probability import FreqDist
from nltk.tokenize import RegexpTokenizer, sent_tokenize
from collections import Counter, defaultdict
from textblob import TextBlob

# Download necessary NLTK resources if not already present
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)
    
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords', quiet=True)
    
try:
    nltk.data.find('taggers/averaged_perceptron_tagger')
except LookupError:
    nltk.download('averaged_perceptron_tagger', quiet=True)
    
try:
    nltk.data.find('corpora/wordnet')
except LookupError:
    nltk.download('wordnet', quiet=True)
    
try:
    nltk.data.find('chunkers/maxent_ne_chunker')
except LookupError:
    nltk.download('maxent_ne_chunker', quiet=True)
    
try:
    nltk.data.find('corpora/words')
except LookupError:
    nltk.download('words', quiet=True)

# Logging with Rich
from rich.console import Console
from rich.logging import RichHandler

# Setup logging
logging.basicConfig(
    level=logging.ERROR,
    format="%(message)s",
    datefmt="[%X]",
    handlers=[RichHandler(rich_tracebacks=True)]
)
logger = logging.getLogger("textextract")

console = Console()

# Use rich for pretty printing by default
print = console.print
log = logger.info

# User agent strings for web requests


def generate_http_headers(url):
    USER_AGENTS = [
        # Desktop
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 13_3) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Safari/605.1.15",
        "Mozilla/5.0 (X11; Linux x86_64) Gecko/20100101 Firefox/115.0",
        # Mobile
        "Mozilla/5.0 (iPhone; CPU iPhone OS 16_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Mobile/15E148 Safari/604.1",
        "Mozilla/5.0 (Linux; Android 13; Pixel 7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.6367.91 Mobile Safari/537.36"
    ]

    headers = {
        "User-Agent": random.choice(USER_AGENTS),
        "Accept-Language": "en-US,en;q=0.9",
        "Referer": url,
        "DNT": "1",
        "Upgrade-Insecure-Requests": "1"
    }

    return headers



def clean_path(
    path: str
) -> Optional[str]:
    """
    Normalize and validate a filesystem path.

    Args:
        path (str): Input file or directory path.

    Returns:
        Optional[str]: Absolute path if valid; None otherwise.
    """
    if is_url(path):
        return path

    p = os.path.expanduser(path)
    p = os.path.abspath(p)
    if os.path.isfile(p) or os.path.isdir(p):
        return p
    return None


def normalize_text(
    text: str
) -> str:
    """
    Replace multiple consecutive newlines, carriage returns, and spaces
    with a single space. Ensures compact, single-line output.

    Args:
        text (str): Raw input text.

    Returns:
        str: Normalized single-line text.
    """
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r' +', ' ', text)
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r'(?m)(^ \n)+', '\n', text)
    text = re.sub(r'\t+', '\t', text)
    text = re.sub(r'\r+', '\n', text)
    text = re.sub(r"^ ", "", text, flags=re.MULTILINE)
    return text 


def is_url(s: str) -> bool:
    """
    Check if a string is a valid HTTP/HTTPS URL.
    
    Args:
        s (str): String to check
        
    Returns:
        bool: True if valid URL, False otherwise
    """
    parsed = urlparse(s)
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def text_from_screenshot() -> str:
    """
    Capture a full-screen screenshot, perform OCR, and clean up temp file.

    Returns:
        str: Normalized OCR-extracted text from the screenshot.
    """
    tmp_filename = f"screenshot_{uuid4().hex}.png"
    tmp_path = os.path.join(tempfile.gettempdir(), tmp_filename)

    try:
        with mss() as sct:
            monitor = {"top": 0, "left": 0, "width": 0, "height": 0}
            for mon in sct.monitors:
                monitor["left"] = min(mon["left"], monitor["left"])
                monitor["top"] = min(mon["top"], monitor["top"])
                monitor["width"] = max(mon["width"] + mon["left"] - monitor["left"], monitor["width"])
                monitor["height"] = max(mon["height"] + mon["top"] - monitor["top"], monitor["height"])
            screenshot = sct.grab(monitor)

        img = Image.frombytes("RGB", screenshot.size, screenshot.bgra, "raw", "BGRX")
        img_gray = img.convert("L")
        img_gray.save(tmp_path)

        content = text_from_image(tmp_path)
        return normalize_text(content)
    finally:
        if os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except Exception as e:
                logger.error(f"Failed to delete temp screenshot: {e}")


def extract_exif(
    file_path: str
) -> Optional[Dict[str, Any]]:
    """
    Extract EXIF metadata from a file using exiftool.

    Args:
        file_path (str): Path to the target file.

    Returns:
        Optional[Dict[str, Any]]: Parsed EXIF data, or None on failure.
    """
    exif_data: Optional[Dict[str, Any]] = None
    try:
        result = subprocess.run(
            ['exiftool', '-j', file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        if result.returncode == 0:
            exif_data = json.loads(result.stdout.decode())[0]
    except Exception as e:
        logger.error(f"Exiftool failed: {e}")
    return exif_data


def extract_document_structure(html: str) -> Dict[str, List[Dict[str, Any]]]:
    """
    Extract document structure (headings, lists, sections) from HTML.
    
    Args:
        html (str): HTML content
        
    Returns:
        Dict: Document structure information
    """
    soup = BeautifulSoup(html, "html.parser")
    
    structure = {
        "headings": [],
        "lists": [],
        "sections": []
    }
    
    # Extract headings
    for level in range(1, 7):
        for heading in soup.find_all(f'h{level}'):
            structure["headings"].append({
                "level": level,
                "text": heading.get_text(strip=True),
                "id": heading.get("id", "")
            })
    
    # Extract lists
    for list_type in ['ul', 'ol']:
        for lst in soup.find_all(list_type):
            items = [li.get_text(strip=True) for li in lst.find_all('li')]
            structure["lists"].append({
                "type": "unordered" if list_type == "ul" else "ordered",
                "items": items
            })
    
    # Extract sections
    for section in soup.find_all('section'):
        section_info = {
            "id": section.get("id", ""),
            "class": section.get("class", []),
            "text": section.get_text(strip=True)
        }
        structure["sections"].append(section_info)
    
    return structure


def text_from_html(html: str) -> str:
    """
    Extract readable text from raw HTML content.

    Args:
        html (str): HTML source as a string.

    Returns:
        str: Cleaned and normalized visible text.
    """
    # Check if the input is a file path or HTML content
    if os.path.isfile(html):
        with open(html, 'r', encoding='utf-8', errors='ignore') as f:
            html = f.read()
    
    soup = BeautifulSoup(html, "html.parser")

    # Remove non-visible or structural elements
    for tag in soup([
        "script", "style",
        "noscript", "iframe",
        "meta", "link",
        "header", "footer",
        "form", "nav",
        "aside"
    ]):
        tag.decompose()

    text = soup.get_text(separator=" ")

    return normalize_text(text)


def text_from_url(
    url: str,
    render_js: bool = True
) -> Optional[str]:
    """
    Extract visible text from a web page or downloadable file at the given URL.

    Args:
        url (str): Target URL (web page or file).
        render_js (bool): Whether to render JavaScript content.

    Returns:
        Optional[str]: Extracted text, or None on failure.
    """
    headers = generate_http_headers(url)

    # Attempt to detect content type
    content_type = ""
    try:
        head = requests.head(url, headers=headers, timeout=5, allow_redirects=True)
        content_type = head.headers.get("Content-Type", "").lower()
    except Exception as e:
        logger.warning(f"HEAD request failed: {e}")
        try:
            resp = requests.get(url, headers=headers, timeout=5, stream=True)
            content_type = resp.headers.get("Content-Type", "").lower()
        except Exception as e:
            logger.warning(f"GET fallback for Content-Type check failed: {e}")

    # If clearly not HTML, treat as a file and extract locally
    if not content_type.startswith("text/html"):
        try:
            with requests.get(url, headers=headers, stream=True, timeout=15) as r:
                r.raise_for_status()
                suffix = Path(urlparse(url).path).suffix or ".bin"
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix, dir="/tmp") as tmp_file:
                    shutil.copyfileobj(r.raw, tmp_file)
                    tmp_path = tmp_file.name
            return extract_text(tmp_path)
        except Exception as e:
            logger.error(f"Failed to download and extract file from URL: {url} - {e}")
            return None

    # Standard HTML path
    if render_js:
        try:
            session = HTMLSession()
            try:
                r = session.get(url, headers=headers, timeout=5)
                try:
                    r.html.render(timeout=5, sleep=1, keep_page=True)
                except Exception as e:
                    logger.warning(f"JS rendering failed, falling back to static HTML: {e}")
                html = r.html.html
                return text_from_html(html)
            except Exception as e:
                logger.error(f"[Error with HTMLSession] {url} - {e}")
            finally:
                session.close()
        except Exception as e:
            logger.error(f"[Error creating HTMLSession] {e}")

    # Fallback: static HTML without rendering
    try:
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        html = response.text
        return text_from_html(html)
    except Exception as e:
        logger.error(f"[Error with requests] {url} - {e}")
        return None



def scrape_website(url: str, max_pages: int = 1, stay_on_domain: bool = True) -> Dict[str, str]:
    """
    Scrape multiple pages of a website.
    
    Args:
        url (str): Starting URL
        max_pages (int): Maximum pages to scrape
        stay_on_domain (bool): Whether to stay on the same domain
        
    Returns:
        Dict[str, str]: Dictionary mapping URLs to extracted text
    """
    results = {}
    visited = set()
    to_visit = [url]
    base_domain = urlparse(url).netloc
    
    while to_visit and len(visited) < max_pages:
        current_url = to_visit.pop(0)
        if current_url in visited:
            continue
            
        # Extract text from current page
        text = text_from_url(current_url)
        if text:
            results[current_url] = text
            
        visited.add(current_url)
        
        # Find links on the page
        session = HTMLSession()
        try:
            r = session.get(current_url)
            r.html.render(timeout=20, sleep=1)
            
            links = r.html.absolute_links
            for link in links:
                link_domain = urlparse(link).netloc
                if link not in visited and link not in to_visit:
                    # Check if we should follow this link
                    if stay_on_domain and link_domain != base_domain:
                        continue
                    to_visit.append(link)
        except Exception as e:
            logger.error(f"Error scraping {current_url}: {e}")
        finally:
            session.close()
    
    return results


def extract_text(
    file_path: str
) -> Optional[str]:
    """
    Extract text content from a local file or URL.

    Supports web pages, text, JSON, XML, CSV, Excel, PDF, DOCX, images, audio.

    Args:
        file_path (str): Path to the input file or URL.

    Returns:
        Optional[str]: Extracted text, or None if unsupported or error.
    """
    if is_url(file_path):
        return text_from_url(file_path)

    TEXT_MIME_TYPES = {
        "application/json", "application/xml", "application/x-yaml",
        "application/x-toml", "application/x-csv", "application/x-markdown",
    }

    path = clean_path(file_path)
    if not path:
        logger.error(f"No such file: {file_path}")
        return None

    mime_type = magic.from_file(path, mime=True)
    try:
        if mime_type.startswith("text/html"):
            content = text_from_html(path)
            return content

        elif mime_type.startswith("text/") or mime_type in TEXT_MIME_TYPES:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return normalize_text(content)

        elif mime_type in [
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ]:
            content = text_from_excel(path)
            return content

        elif mime_type == "application/pdf":
            content = text_from_pdf(path)
            return content

        elif mime_type == \
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            content = text_from_docx(path)
            return content

        elif mime_type == "application/msword":
            content = text_from_doc(path)
            return content

        elif mime_type.startswith("image/"):
            content = text_from_image(path)
            return content

        elif mime_type.startswith("audio/"):
            content = text_from_audio(path)
            return content

        elif mime_type == "application/epub+zip":
            content = text_from_epub(path)
            return content

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            content = text_from_pptx(path)
            return content

        elif mime_type == "application/vnd.oasis.opendocument.text":
            content = text_from_odt(path)
            return content

        else:
            content = text_from_any(path)
            return content
    except Exception as e:
        logger.error(f"Error reading {path}: {e}")
        return None


def extract_text_with_password(file_path: str, password: str) -> Optional[str]:
    """
    Extract text from password-protected files.
    
    Args:
        file_path (str): Path to the file
        password (str): Password to unlock the file
        
    Returns:
        Optional[str]: Extracted text
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.pdf':
        return text_from_pdf_protected(file_path, password)
    elif file_ext in ['.docx', '.xlsx', '.pptx']:
        return text_from_office_protected(file_path, password)
    else:
        logger.warning(f"Password protection not supported for {file_ext} files")
        return None


def text_from_pdf_protected(pdf_path: str, password: str) -> Optional[str]:
    """
    Extract text from password-protected PDF.
    
    Args:
        pdf_path (str): Path to the PDF file
        password (str): Password to unlock the PDF
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        doc = pymupdf.open(pdf_path, password=password)
        text = ""
        for i in range(len(doc)):
            page = doc.load_page(i)
            text += page.get_text()
        doc.close()
        return normalize_text(text)
    except Exception as e:
        logger.error(f"Error in PDF with password: {str(e)}")
        return None


def text_from_office_protected(office_path: str, password: str) -> Optional[str]:
    """
    Extract text from password-protected Office files.
    
    Args:
        office_path (str): Path to the Office file
        password (str): Password to unlock the file
        
    Returns:
        Optional[str]: Extracted text
    """
    # This is a placeholder. Actual implementation would depend on the specific office format
    # and available libraries for password-protected office files.
    logger.warning("Password-protected Office files extraction not fully implemented")
    return None


def text_from_audio(
    audio_file: str
) -> Optional[str]:
    """
    Transcribe audio to text using Google Speech Recognition.

    Args:
        audio_file (str): Path to the input audio file.

    Returns:
        Optional[str]: Transcription, or None on failure.
    """
    def convert_to_wav(file_path: str) -> str:
        _, ext = os.path.splitext(file_path)
        ext = ext.lstrip('.')
        audio = AudioSegment.from_file(file_path, format=ext)
        tmp_filename = f"audio_{uuid4().hex}.wav"
        wav_path = os.path.join(tempfile.gettempdir(), tmp_filename)
        audio.export(wav_path, format='wav')
        return wav_path

    recognizer = sr.Recognizer()
    temp_wav_path = None
    cleanup_needed = False

    try:
        _, ext = os.path.splitext(audio_file)
        if ext.lower() not in ['.wav', '.wave']:
            temp_wav_path = convert_to_wav(audio_file)
            cleanup_needed = True
        else:
            temp_wav_path = clean_path(audio_file)

        if not temp_wav_path:
            logger.error("Invalid audio path.")
            return None

        with sr.AudioFile(temp_wav_path) as source:
            audio = recognizer.record(source)
        return recognizer.recognize_google(audio)

    except sr.UnknownValueError:
        logger.error("Could not understand audio.")
    except sr.RequestError as e:
        logger.error(f"Speech recognition error: {e}")
    except Exception as e:
        logger.error(f"Failed to process audio: {e}")
    finally:
        if cleanup_needed and temp_wav_path and os.path.exists(temp_wav_path):
            try:
                os.remove(temp_wav_path)
            except Exception as e:
                logger.error(f"Failed to delete temp WAV file {temp_wav_path}: {e}")

    return None


def downloadImage(
    url: str
) -> Optional[str]:
    """
    Download an image from a URL to /tmp/ and return its path.

    Args:
        url (str): Remote image URL.

    Returns:
        Optional[str]: Local file path, or None on failure.
    """
    if is_url(url):
        try:
            resp = requests.head(url)
            content_type = resp.headers.get('Content-Type', '')
            if content_type.startswith('image/'):
                filename = os.path.basename(urlparse(url).path)
                if not filename:
                    filename = f"image_{uuid4().hex}.jpg"  # Default name if none found
                save_path = os.path.join(tempfile.gettempdir(), filename)
                resp = requests.get(url, stream=True)
                resp.raise_for_status()
                with open(save_path, "wb") as f:
                    for chunk in resp.iter_content(chunk_size=8192):
                        f.write(chunk)
                return clean_path(save_path)
        except Exception as e:
            logger.error(f"Failed to download image from {url}: {e}")
    logger.error(f"Unable to pull image from {url}")
    return None


def is_image(
    file_path_or_url: str
) -> bool:
    """
    Determine if the given path/URL points to an image.

    Args:
        file_path_or_url (str): Local path or URL.

    Returns:
        bool: True if MIME type starts with 'image/'.
    """
    try:
        if is_url(file_path_or_url):
            resp = requests.head(file_path_or_url)
            content_type = resp.headers.get('Content-Type', '')
            return content_type.startswith('image/')
        else:
            mime = magic.from_file(file_path_or_url, mime=True)
            return mime.startswith("image/")
    except Exception:
        return False


def text_from_pdf(
    pdf_path: str
) -> Optional[str]:
    """
    Extract text and OCR results from a PDF using PyMuPDF.

    Args:
        pdf_path (str): Path to PDF file.

    Returns:
        Optional[str]: Combined normalized text and image OCR results.
    """
    plain_text = ""
    temp_image_paths: List[str] = []

    try:
        doc = pymupdf.open(pdf_path)
        for k, v in doc.metadata.items():
            plain_text += f"{k}: {v}\n"

        for i in range(len(doc)):
            page = doc.load_page(i)
            plain_text += f"\n--- Page {i + 1} ---\n"
            text = page.get_text()
            plain_text += text or "[No text]\n"

            for img_index, img in enumerate(page.get_images(full=True), start=1):
                xref = img[0]
                base = doc.extract_image(xref)
                img_bytes = base["image"]

                img_filename = f"pdf_page{i+1}_img{img_index}_{uuid4().hex}.png"
                img_path = os.path.join(tempfile.gettempdir(), img_filename)
                temp_image_paths.append(img_path)

                with open(img_path, "wb") as f:
                    f.write(img_bytes)

                ocr = text_from_image(img_path) or ""
                plain_text += f"\n[Image {img_index} OCR]\n{ocr}\n"

        # Extract tables from PDF
        """
        try:
            tables = extract_tables_from_pdf(pdf_path)
            if tables:
                plain_text += "\n--- Tables ---\n"
                for i, table in enumerate(tables, 1):
                    plain_text += f"\n[Table {i}]\n"
                    if isinstance(table, dict) and "data" in table:
                        for row in table["data"]:
                            plain_text += str(row) + "\n"
                    else:
                        plain_text += str(table) + "\n"
        except Exception as e:
            logger.warning(f"Could not extract tables from PDF: {e}")
        """

        return normalize_text(plain_text)
    except Exception as e:
        logger.error(f"Error processing PDF: {e}")
        return None
    finally:
        for path in temp_image_paths:
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception as e:
                    logger.error(f"Failed to delete temp image {path}: {e}")
        if 'doc' in locals():
            doc.close()


def extract_tables_from_pdf(pdf_path: str) -> List[Dict[str, Any]]:
    """
    Extract tables from PDF using tabula-py.
    
    Args:
        pdf_path (str): Path to the PDF file
        
    Returns:
        List[Dict[str, Any]]: List of extracted tables
    """
    try:
        import tabula
        
        tables = []
        # Read all tables from the PDF
        extracted_tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
        
        for i, table in enumerate(extracted_tables):
            if not table.empty:
                tables.append({
                    "table_number": i+1,
                    "rows": len(table),
                    "columns": len(table.columns),
                    "data": table.to_dict(orient='records')
                })
        
        return tables
    except ImportError:
        logger.warning("tabula-py not installed, skipping table extraction")
        return []
    except Exception as e:
        logger.error(f"Error extracting tables from PDF: {e}")
        return []


def extract_pdf_chunked(pdf_path: str, chunk_size: int = 10) -> Optional[str]:
    """
    Extract text from large PDFs in chunks to manage memory usage.
    
    Args:
        pdf_path (str): Path to the PDF file
        chunk_size (int): Number of pages to process at once
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        doc = pymupdf.open(pdf_path)
        total_pages = len(doc)
        text_chunks = []
        
        # Add metadata
        metadata = ""
        for k, v in doc.metadata.items():
            metadata += f"{k}: {v}\n"
        text_chunks.append(metadata)
        
        for i in range(0, total_pages, chunk_size):
            chunk_text = ""
            end_page = min(i + chunk_size, total_pages)
            
            for page_num in range(i, end_page):
                page = doc.load_page(page_num)
                chunk_text += f"\n--- Page {page_num + 1} ---\n"
                chunk_text += page.get_text() or "[No text]\n"
                
            text_chunks.append(chunk_text)
            
        doc.close()
        return normalize_text("".join(text_chunks))
    except Exception as e:
        logger.error(f"Error processing PDF in chunks: {e}")
        return None


def text_from_doc(
    filepath: str,
    min_length: int = 4
) -> str:
    """
    Extract readable strings and metadata from binary Word (.doc) files.

    Args:
        filepath (str): Path to .doc file.
        min_length (int): Minimum string length to extract.

    Returns:
        str: Metadata and text content.
    """
    def extract_printable_strings(
        data: bytes
    ) -> List[str]:
        pattern = re.compile(
            b'[' + re.escape(bytes(string.printable, 'ascii')) +
            b']{%d,}' % min_length
        )
        found = pattern.findall(data)

        results = []
        for m in found:
            value = m.decode(errors='ignore').strip()
            results.append(value)

        return results

    def clean_strings(
        strs: List[str]
    ) -> List[str]:
        cleaned: List[str] = []
        skip = ["HYPERLINK", "OLE2", "Normal.dotm"]
        for line in strs:
            if any(line.startswith(pref) for pref in skip):
                continue
            cleaned.append(re.sub(r'\s+', ' ', line).strip())
        return cleaned

    with open(filepath, 'rb') as f:
        data = f.read()

    strings = extract_printable_strings(data)
    strings = clean_strings(strings)
    content = "\n".join(strings)

    return normalize_text(content)


def text_from_docx(
    file_path: str
) -> Optional[str]:
    """
    Extract text, tables, and OCR from embedded images in a DOCX file.

    Args:
        file_path (str): Path to the .docx file.

    Returns:
        Optional[str]: Normalized full text content.
    """
    path = clean_path(file_path)
    if not path:
        return None

    temp_image_paths: List[str] = []
    plain_text = ""

    try:
        doc = Document(path)

        for p in doc.paragraphs:
            if p.text.strip():
                plain_text += p.text.strip() + "\n"

        for tbl in doc.tables:
            plain_text += "\n[Table]\n"
            for row in tbl.rows:
                row_text = "\t".join(c.text.strip() for c in row.cells)
                plain_text += row_text + "\n"

        for rel_id, rel in doc.part.rels.items():
            if "image" in rel.target_ref:
                blob = rel.target_part.blob

                img_filename = f"docx_img_{rel_id}_{uuid4().hex}.png"
                img_path = os.path.join(tempfile.gettempdir(), img_filename)
                temp_image_paths.append(img_path)

                with open(img_path, "wb") as img_file:
                    img_file.write(blob)

                ocr = text_from_image(img_path) or ""
                plain_text += f"\n[Image OCR]\n{ocr}\n"

        return normalize_text(plain_text)

    except Exception as e:
        logger.error(f"Error processing DOCX: {e}")
        return None
    finally:
        for path in temp_image_paths:
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception as e:
                    logger.error(f"Failed to delete temp DOCX image {path}: {e}")


def text_from_excel(
    file_path: str
) -> str:
    """
    Convert an Excel workbook to CSV text.

    Args:
        file_path (str): Path to the Excel file.

    Returns:
        str: CSV-formatted string.
    """
    path = clean_path(file_path)
    if not path:
        return ""
    try:
        # Get all sheets
        result = ""
        excel_file = pd.ExcelFile(path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet_name)
            out = StringIO()
            df.to_csv(out, index=False)
            result += f"\n--- Sheet: {sheet_name} ---\n"
            result += out.getvalue()
            result += "\n"
        return result
    except Exception as e:
        logger.error(f"Failed Excel -> CSV: {e}")
        return ""


def text_from_image(
    file_path: str
) -> Optional[str]:
    """
    Perform OCR on an image file.

    Args:
        file_path (str): Path to the image.

    Returns:
        Optional[str]: Extracted text, or None on error.
    """
    path = clean_path(file_path)
    if not path:
        return None
    try:
        with Image.open(path) as img:
            # Improve OCR with preprocessing
            # 1. Convert to grayscale if it's not already
            if img.mode != 'L':
                img = img.convert('L')
                
            # 2. Optional: Apply some contrast enhancement
            # (Disabled by default, enable if needed for specific cases)
            # from PIL import ImageEnhance
            # enhancer = ImageEnhance.Contrast(img)
            # img = enhancer.enhance(1.5)  # Increase contrast
                
            # Perform OCR with custom configuration
            custom_config = r'--oem 3 --psm 6'  # Default OCR Engine Mode and Page Segmentation Mode
            txt = pytesseract.image_to_string(img, config=custom_config).strip()
            return normalize_text(txt) or ""
    except Exception as e:
        logger.error(f"Failed image OCR: {e}")
        return None


def extract_strings(file_path, min_length=4):
    """
    Extract printable strings from a file, similar to the Unix 'strings' command.
    
    Args:
        file_path (str): Path to the file to extract strings from
        min_length (int, optional): Minimum length of strings to extract. Defaults to 4.
        
    Returns:
        list: List of printable strings found in the file
    """
    file_path = clean_path(file_path)

    
    # Define printable characters (excluding tabs and newlines)
    printable_chars = set(string.printable) - set('\t\n\r\v\f')
    
    result = []
    current_string = ""
    
    # Read the file in binary mode
    try:
        with open(file_path, 'rb') as file:
            # Read the file byte by byte
            for byte in file.read():
                # Convert byte to character
                char = chr(byte)
                
                # If character is printable, add to current string
                if char in printable_chars:
                    current_string += char
                # If not printable and we have a string of minimum length, add to results
                elif len(current_string) >= min_length:
                    if current_string == "Sj[d":
                        pass
                    else:
                        result.append(current_string)
                    current_string = ""
                # If not printable and current string is too short, reset current string
                else:
                    current_string = ""
        
        # Don't forget to add the last string if it meets the minimum length
        if len(current_string) >= min_length:
            result.append(current_string)
        
        return result
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.", file=sys.stderr)
        return None
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        return None


def text_from_any(
    file_path: str
) -> Optional[str]:
    """
    Handle unknown file types by reporting stats and metadata.

    Args:
        file_path (str): Path to the file.

    Returns:
        Optional[str]: Plain-text report, or None on error.
    """
    content = ""
    path = clean_path(file_path)
    if not path:
        return None
    try:
        stats = os.stat(path)
        info = {
            "path": path,
            "size": stats.st_size,
            "created": datetime.fromtimestamp(stats.st_ctime).isoformat(),
            "modified": datetime.fromtimestamp(stats.st_mtime).isoformat(),
        }

        for k, v in info.items():
            content += "File System Data:\n"
            content += f"{k}: {v}\n"
        
        # Try to extract EXIF if available
        exif = extract_exif(path)
        if exif:
            info["exif"] = exif
            content += "\n\nEXIF Data:\n"
            for k, v in exif.items():
                if isinstance(v, dict):
                    content += f"\n{k}:\n"
                    for sub_k, sub_v in v.items():
                        content += f"  {sub_k}: {sub_v}\n"
                else:
                    content += f"{k}: {v}\n"

        # Get file hash
        md5_hash = hashlib.md5(open(path,'rb').read()).hexdigest()
        info["md5"] = md5_hash

        # Get strings
        strings = extract_strings(path)
        if strings:
            info["strings"] = strings
            content += "\n\nStrings Data:\n"
            clean_strings = "\n".join(strings)
            content += clean_strings

        return info 
    except Exception as e:
        logger.error(f"Error on other file: {e}")
        return None


def text_from_object(obj):
    return json.dumps(obj, indent=4, ensure_ascii=False)

def tree_from_object(obj, indent=0, path=None):
    """
    Format a Python object (dict, list, tuple, set) in a visually appealing,
    human-readable way with excellent handling of nested structures.
    
    Args:
        obj: The Python object to format
        indent: Current indentation level
        path: Current path in the object (for tracking circular references)
        
    Returns:
        String representation of the object formatted for human readability
    """
    def _generate_object_summary(obj):
        """Generate a summary of the object's structure and contents"""
        # Initialize summary sections
        summary_lines = []
        summary_lines.append(" OBJECT SUMMARY")
        summary_lines.append("╔═══════════════")
        
        # Get basic type information
        obj_type = type(obj).__name__
        
        # Object type and size
        if isinstance(obj, (dict, list, tuple, set)):
            summary_lines.append(f"║ Type:           {obj_type:<60}")
            summary_lines.append(f"║ Item Count:     {len(obj):<60}")
        else:
            summary_lines.append(f"║ Type:           {obj_type:<60}")
            summary_lines.append(f"║ Size:           {len(str(obj)) if hasattr(obj, '__len__') else 'N/A':<60}")
        
        # Get structure info for collections
        if isinstance(obj, dict):
            # Analyze keys
            key_types = {}
            for k in obj.keys():
                k_type = type(k).__name__
                key_types[k_type] = key_types.get(k_type, 0) + 1
            
            key_type_str = ", ".join(f"{k}: {v}" for k, v in key_types.items())
            summary_lines.append(f"║ Key Types:      {key_type_str:<60}")
            
            # Analyze values
            value_types = {}
            nested_counts = {'dict': 0, 'list': 0, 'tuple': 0, 'set': 0}
            
            for v in obj.values():
                v_type = type(v).__name__
                value_types[v_type] = value_types.get(v_type, 0) + 1
                
                # Count nested collections
                if v_type in nested_counts:
                    nested_counts[v_type] += 1
            
            value_type_str = ", ".join(f"{k}: {v}" for k, v in value_types.items())
            if len(value_type_str) > 28:
                value_type_str = value_type_str[:38] + "..."
            summary_lines.append(f"║ Value Types:    {value_type_str:<60}")
            
            # Show nesting info
            nested_str = ", ".join(f"{k}: {v}" for k, v in nested_counts.items() if v > 0)
            if nested_str:
                summary_lines.append(f"║ Nested Objects: {nested_str:<60}")
        
        elif isinstance(obj, (list, tuple)):
            # Analyze item types
            item_types = {}
            nested_counts = {'dict': 0, 'list': 0, 'tuple': 0, 'set': 0}
            
            for item in obj:
                item_type = type(item).__name__
                item_types[item_type] = item_types.get(item_type, 0) + 1
                
                # Count nested collections
                if item_type in nested_counts:
                    nested_counts[item_type] += 1
            
            item_type_str = ", ".join(f"{k}: {v}" for k, v in item_types.items())
            if len(item_type_str) > 28:
                item_type_str = item_type_str[:38] + "..."
            summary_lines.append(f"║ Item Types:     {item_type_str:<60}")
            
            # Show nesting info
            nested_str = ", ".join(f"{k}: {v}" for k, v in nested_counts.items() if v > 0)
            if nested_str:
                summary_lines.append(f"║ Nested Objects: {nested_str:<60}")
        
        elif isinstance(obj, set):
            # Analyze item types
            item_types = {}
            nested_counts = {'dict': 0, 'list': 0, 'tuple': 0, 'set': 0}
            
            for item in obj:
                item_type = type(item).__name__
                item_types[item_type] = item_types.get(item_type, 0) + 1
                
                # Count nested collections
                if item_type in nested_counts:
                    nested_counts[item_type] += 1
            
            item_type_str = ", ".join(f"{k}: {v}" for k, v in item_types.items())
            if len(item_type_str) > 28:
                item_type_str = item_type_str[:38] + "..."
            summary_lines.append(f"║ Item Types:     {item_type_str:<60}")
            
            # Show nesting info
            nested_str = ", ".join(f"{k}: {v}" for k, v in nested_counts.items() if v > 0)
            if nested_str:
                summary_lines.append(f"║ Nested Objects: {nested_str:<60}")
        
        # Add max nesting depth for collections
        if isinstance(obj, (dict, list, tuple, set)):
            max_depth = _calculate_max_depth(obj)
            summary_lines.append(f"║ Max Nest Depth: {max_depth:<60}")
        
        # Close the summary box
        summary_lines.append("║")
        
        return "\n".join(summary_lines)

    def _calculate_max_depth(obj, current_depth=1, visited=None):
        """Calculate the maximum nesting depth of a collection"""
        if visited is None:
            visited = set()
        
        # Handle circular references
        obj_id = id(obj)
        if obj_id in visited:
            return current_depth
        
        if not isinstance(obj, (dict, list, tuple, set)):
            return current_depth
        
        # Add to visited set
        visited.add(obj_id)
        
        # Check depth of nested objects
        max_nested_depth = current_depth
        
        if isinstance(obj, dict):
            for value in obj.values():
                if isinstance(value, (dict, list, tuple, set)):
                    depth = _calculate_max_depth(value, current_depth + 1, visited)
                    max_nested_depth = max(max_nested_depth, depth)
        
        elif isinstance(obj, (list, tuple, set)):
            for item in obj:
                if isinstance(item, (dict, list, tuple, set)):
                    depth = _calculate_max_depth(item, current_depth + 1, visited)
                    max_nested_depth = max(max_nested_depth, depth)
        
        return max_nested_depth

    def _format_object_detail(obj, indent=0, path=None):
        """Format the detailed view of the object"""
        if path is None:
            path = []
        
        # Characters for tree structure
        indent_str = " " * indent
        first_branch = "┌─ "  # First item branch
        branch = "├─ "        # Middle item branch
        last_branch = "└─ "   # Last item branch
        vertical = "│  "      # Vertical connector
        space = " "         # Space (no vertical)
        
        # Prevent infinite recursion with circular references
        obj_id = id(obj)
        if obj_id in path and isinstance(obj, (dict, list, tuple, set)):
            return f"{indent_str}<circular reference>"
        
        # Create new path for recursive calls
        new_path = path + [obj_id]
        result = []
        
        # Format None
        if obj is None:
            return "None"
        
        # Format primitive types
        if isinstance(obj, (int, float, bool, str)):
            if isinstance(obj, str):
                if '\n' in obj:
                    # Format multiline strings
                    lines = obj.split('\n')
                    if len(lines) > 3:
                        preview = f"{lines[0]}... [+{len(lines)-1} more lines]"
                        return f'"{preview}"'
                    return f'"""{obj}"""'
                return f'"{obj}"'
            elif isinstance(obj, bool):
                return "True" if obj else "False"
            elif isinstance(obj, int) and abs(obj) >= 1000:
                return f"{obj:,}"
            elif isinstance(obj, float):
                if abs(obj) < 0.001 or abs(obj) >= 10000:
                    return f"{obj:.6e}"
                return f"{obj:.6g}"
            return str(obj)
        
        # Format dictionaries
        if isinstance(obj, dict):
            if not obj:
                return "{}"
            
            # Add type header for top-level objects
            if indent == 0:
                result.append(f"Dictionary ({len(obj)} items)")
                result.append("┌────────────────────")
            
            items = list(obj.items())
            
            # Find the longest key for alignment (if keys are strings)
            key_str_lengths = [len(str(k)) for k in obj.keys()]
            if key_str_lengths:
                max_key_length = max(key_str_lengths)
            else:
                max_key_length = 0
            
            # Process all items
            for i, (key, value) in enumerate(items):
                # Determine the appropriate branch character
                if i == 0 and i == len(items) - 1:  # Only item
                    prefix = "──"  # Single item
                    next_prefix = space
                elif i == 0:  # First item
                    prefix = first_branch 
                    next_prefix = vertical
                elif i == len(items) - 1:  # Last item
                    prefix = last_branch
                    next_prefix = space
                else:  # Middle items
                    prefix = branch
                    next_prefix = vertical
                
                # Format key
                key_str = _format_primitive_for_key(key)
                
                # Handle different value types
                if isinstance(value, (dict, list, tuple, set)):
                    if not value:  # Empty container
                        type_name = type(value).__name__
                        empty_repr = {"dict": "{}", "list": "[]", "tuple": "()", "set": "set()"}
                        result.append(f"{indent_str}{prefix}{key_str}:")
                    else:
                        # Add key with container type and count
                        type_name = type(value).__name__
                        count = len(value)
                        result.append(f"{indent_str}{prefix}{key_str}:") 

                        # Add nested container elements with proper indentation
                        nested_lines = _format_object_detail(value, indent + 1, new_path).split('\n')
                        for line in nested_lines:
                            result.append(f"{indent_str}{next_prefix}{line}")
                else:
                    # Format primitive value with alignment
                    value_str = _format_primitive_for_value(value)
                    padded_key = f"{key_str}:"
                    if all(isinstance(k, str) for k in obj.keys()):
                        padded_key = f"{key_str}:".ljust(max_key_length + 2)
                    result.append(f"{indent_str}{prefix}{padded_key} {value_str}")
            
            return '\n'.join(result)
        
        # Format lists
        if isinstance(obj, list):
            if not obj:
                return "[]"
            
            # Add type header for top-level objects
            if indent == 0:
                result.append(f"List ({len(obj)} items)")
                result.append("┌────────────────────")
            
            # If list contains only primitives and is small, format compactly
            all_primitives = all(not isinstance(x, (dict, list, tuple, set)) for x in obj)
            if all_primitives and len(obj) <= 7 and all(len(str(_format_primitive_for_value(x))) < 15 for x in obj):
                items_str = ", ".join(_format_primitive_for_value(x) for x in obj)
                if indent == 0:
                    return f"List ({len(obj)} items)\n┌────────────────────\n{indent_str}[{items_str}]"
                else:
                    return f"[{items_str}]"
            
            # Otherwise, format each item on its own line
            for i, item in enumerate(obj):
                # Determine the appropriate branch character
                if i == 0 and i == len(obj) - 1:  # Only item
                    prefix = "──"  # Single item
                    next_prefix = space
                elif i == 0:  # First item
                    prefix = first_branch
                    next_prefix = vertical
                elif i == len(obj) - 1:  # Last item
                    prefix = last_branch
                    next_prefix = space
                else:  # Middle items
                    prefix = branch
                    next_prefix = vertical
                
                # Add index
                index_str = f"[{i}]"
                
                # Handle different item types
                if isinstance(item, (dict, list, tuple, set)):
                    if not item:  # Empty container
                        type_name = type(item).__name__
                        empty_repr = {"dict": "{}", "list": "[]", "tuple": "()", "set": "set()"}
                        result.append(f"{indent_str}{prefix}{index_str}:")
                    else:
                        # Add index with container type and count
                        type_name = type(item).__name__
                        count = len(item)
                        result.append(f"{indent_str}{prefix}{index_str}:")
                        
                        # Add nested container elements with proper indentation
                        nested_lines = _format_object_detail(item, indent + 1, new_path).split('\n')
                        for line in nested_lines:
                            result.append(f"{indent_str}{next_prefix}{line}")
                else:
                    # Format primitive item
                    value_str = _format_primitive_for_value(item)
                    result.append(f"{indent_str}{prefix}{index_str}: {value_str}")
            
            return '\n'.join(result)
        
        # Format tuples (similar to lists but with different notation)
        if isinstance(obj, tuple):
            if not obj:
                return "()"
            
            # Add type header for top-level objects
            if indent == 0:
                result.append(f"Tuple ({len(obj)} items)")
                result.append("┌────────────────────")
            
            # THE CRITICAL FIX: Always format tuples with tree structure
            # rather than as a single line when they contain primitive types
            #
            # Previously, it had this conditional logic to format small/primitive tuples in a compact way
            # which was inconsistently applied to tuple elements in other structures
            """
            all_primitives = all(not isinstance(x, (dict, list, tuple, set)) for x in obj)
            if all_primitives and len(obj) <= 7 and all(len(str(_format_primitive_for_value(x))) < 15 for x in obj):
                items_str = ", ".join(_format_primitive_for_value(x) for x in obj)
                if indent == 0:
                    return f"Tuple ({len(obj)} items)\n┌────────────────────\n{indent_str}({items_str})"
                else:
                    return f"({items_str})"
            """
            
            # Instead, always format tuples in tree structure with proper indentation
            # Format each item on its own line for consistency
            for i, item in enumerate(obj):
                # Determine the appropriate branch character
                if i == 0 and i == len(obj) - 1:  # Only item
                    prefix = "──"  # Single item
                    next_prefix = space
                elif i == 0:  # First item
                    prefix = first_branch
                    next_prefix = vertical
                elif i == len(obj) - 1:  # Last item
                    prefix = last_branch
                    next_prefix = space
                else:  # Middle items
                    prefix = branch
                    next_prefix = vertical
                
                # Add index
                index_str = f"({i})"
                
                # Handle different item types
                if isinstance(item, (dict, list, tuple, set)):
                    if not item:  # Empty container
                        type_name = type(item).__name__
                        empty_repr = {"dict": "{}", "list": "[]", "tuple": "()", "set": "set()"}
                        result.append(f"{indent_str}{prefix}{index_str}:")
                    else:
                        # Add index with container type and count
                        type_name = type(item).__name__
                        count = len(item)
                        result.append(f"{indent_str}{prefix}{index_str}:")
                        
                        # Add nested container elements with proper indentation
                        nested_lines = _format_object_detail(item, indent + 1, new_path).split('\n')
                        for line in nested_lines:
                            result.append(f"{indent_str}{next_prefix}{line}")
                else:
                    # Format primitive item
                    value_str = _format_primitive_for_value(item)
                    result.append(f"{indent_str}{prefix}{index_str}: {value_str}")
            
            return '\n'.join(result)
        
        # Format sets
        if isinstance(obj, set):
            if not obj:
                return "set()"
            
            # Add type header for top-level objects
            if indent == 0:
                result.append(f"Set ({len(obj)} items)")
                result.append("┌────────────────────")
            
            # If set contains only primitives and is small, format compactly
            all_primitives = all(not isinstance(x, (dict, list, tuple, set)) for x in obj)
            if all_primitives and len(obj) <= 7 and all(len(str(_format_primitive_for_value(x))) < 15 for x in obj):
                # Sort for consistent display
                sorted_items = sorted(obj, key=lambda x: str(x))
                items_str = ", ".join(_format_primitive_for_value(x) for x in sorted_items)
                if indent == 0:
                    return f"Set ({len(obj)} items)\n┌────────────────────\n{indent_str}{{{items_str}}}"
                else:
                    return f"{{{items_str}}}"
            
            # Otherwise, format each item on its own line
            # Sort for consistent display
            items = sorted(list(obj), key=lambda x: str(x))
            
            for i, item in enumerate(items):
                # Determine the appropriate branch character
                if i == 0 and i == len(items) - 1:  # Only item
                    prefix = "──"  # Single item
                    next_prefix = space
                elif i == 0:  # First item
                    prefix = first_branch
                    next_prefix = vertical
                elif i == len(items) - 1:  # Last item
                    prefix = last_branch
                    next_prefix = space
                else:  # Middle items
                    prefix = branch
                    next_prefix = vertical
                
                # Handle different item types
                if isinstance(item, (dict, list, tuple, set)):
                    if not item:  # Empty container
                        type_name = type(item).__name__
                        empty_repr = {"dict": "{}", "list": "[]", "tuple": "()", "set": "set()"}
                        result.append(f"{indent_str}{prefix}{type_name}:")
                    else:
                        # Add container type and count
                        type_name = type(item).__name__
                        count = len(item)
                        result.append(f"{indent_str}{prefix}{type_name}:")
                        
                        # Add nested container elements with proper indentation
                        nested_lines = _format_object_detail(item, indent + 1, new_path).split('\n')
                        for line in nested_lines:
                            result.append(f"{indent_str}{next_prefix}{line}")
                else:
                    # Format primitive item
                    value_str = _format_primitive_for_value(item)
                    result.append(f"{indent_str}{prefix}{value_str}")
            
            return '\n'.join(result)
        
        # Format other types
        return f"{type(obj).__name__}({str(obj)})"

    def _format_primitive_for_key(value):
        """Helper function to format primitive values for dictionary keys"""
        if isinstance(value, str):
            return value
        return str(value)

    def _format_primitive_for_value(value):
        """Helper function to format primitive values"""
        if isinstance(value, str):
            if '\n' in value:
                # Format multiline strings
                lines = value.split('\n')
                if len(lines) > 3:
                    preview = f"{lines[0]}... [+{len(lines)-1} more lines]"
                    return f'"{preview}"'
                return f'"""{value}"""'
            return f'"{value}"'
        elif isinstance(value, bool):
            return "True" if value else "False"
        elif isinstance(value, int) and abs(value) >= 1000:
            return f"{value:,}"
        elif isinstance(value, float):
            if abs(value) < 0.001 or abs(value) >= 10000:
                return f"{value:.6e}"
            return f"{value:.6g}"
        return str(value)

    # Generate summary at top level only
    if indent == 0:
        summary = _generate_object_summary(obj)
        detail = _format_object_detail(obj, indent, path)
        return f"{summary}\n\n{detail}"
    else:
        return _format_object_detail(obj, indent, path)


# Let's define the full function with maximal metadata collection from a URL using requests + BeautifulSoup

def extract_url_metadata(
    url: str,
    timeout: int = 15
) -> dict:
    """
    Collect comprehensive metadata from a given URL using static analysis (no JavaScript rendering).

    This function performs an HTTP GET request to the specified URL, extracts server response data,
    HTML head metadata, structural features, and network/domain information. It supports both standard
    HTML pages and metadata-rich web documents.

    Args:
        url (str): The target URL to inspect.
        timeout (int): Maximum number of seconds to wait for the HTTP response (default: 15).

    Returns:
        dict: A structured dictionary containing:
            - url: The original input URL.
            - timestamp_utc: The ISO timestamp of the request.
            - domain: Parsed domain components and resolved IP address.
            - request: Final URL, method, status code, redirect chain, and latency.
            - response: Server headers including content type, length, date, and encoding.
            - headers: All response headers as a key-value map.
            - html_meta: Parsed metadata from the HTML head, including:
                - <title>, meta[name=...], meta[property=...]
                - Open Graph and Twitter card tags
                - Canonical link, favicon, charset, viewport, robots, hreflangs
                - Language attribute, form presence, media counts, and script count
                - Internal and external hyperlinks
            - network: Reserved for future expansion (e.g., TLS, WHOIS, geolocation)
            - error: Error message if any exception was raised during processing.

    Notes:
        - This does not execute JavaScript or handle SPA/JS-rendered content.
        - Use a headless browser (e.g., Playwright) for dynamic metadata extraction.

    Example:
        metadata = collect_comprehensive_webpage_metadata("https://example.com")
    """
    headers = generate_http_headers(url)

    metadata = {
        "url": url,
        "timestamp_utc": datetime.utcnow().isoformat(),
        "domain": {},
        "request": {},
        "response": {},
        "headers": {},
        "html_meta": {},
        "network": {},
        "error": None
    }

    try:
        # Resolve hostname
        parsed = urlparse(url)
        host = parsed.hostname
        ip = socket.gethostbyname(host) if host else None
        metadata["domain"]["hostname"] = host
        metadata["domain"]["ip_address"] = ip
        metadata["domain"]["scheme"] = parsed.scheme
        metadata["domain"]["path"] = parsed.path
        metadata["domain"]["query"] = parsed.query
        metadata["domain"]["netloc"] = parsed.netloc

        # Start request
        response = requests.get(url, headers=headers, timeout=timeout, allow_redirects=True)

        # Request info
        metadata["request"]["final_url"] = response.url
        metadata["request"]["status_code"] = response.status_code
        metadata["request"]["method"] = response.request.method
        metadata["request"]["redirect_chain"] = [r.url for r in response.history]
        metadata["request"]["elapsed_seconds"] = response.elapsed.total_seconds()

        # Response info
        metadata["response"]["content_type"] = response.headers.get("Content-Type", "")
        metadata["response"]["encoding"] = response.encoding
        metadata["response"]["content_length"] = response.headers.get("Content-Length")
        metadata["response"]["server"] = response.headers.get("Server")
        metadata["response"]["date"] = response.headers.get("Date")

        # Full headers
        metadata["headers"] = dict(response.headers)

        # HTML meta parsing
        soup = BeautifulSoup(response.text, "html.parser")
        head = soup.head or soup

        def get_meta(name=None, prop=None):
            if name:
                tag = head.find("meta", attrs={"name": name})
            elif prop:
                tag = head.find("meta", attrs={"property": prop})
            else:
                return None
            return tag["content"].strip() if tag and tag.has_attr("content") else None

        # Standard meta
        metadata["html_meta"]["title"] = head.title.string.strip() if head.title and head.title.string else None
        metadata["html_meta"]["description"] = get_meta(name="description")
        metadata["html_meta"]["keywords"] = get_meta(name="keywords")
        metadata["html_meta"]["author"] = get_meta(name="author")
        metadata["html_meta"]["robots"] = get_meta(name="robots")
        metadata["html_meta"]["generator"] = get_meta(name="generator")
        metadata["html_meta"]["viewport"] = get_meta(name="viewport")
        metadata["html_meta"]["charset"] = (head.find("meta", charset=True) or {}).get("charset")

        # Open Graph
        for tag in ["og:title", "og:description", "og:image", "og:url", "og:type", "og:site_name"]:
            metadata["html_meta"][tag] = get_meta(prop=tag)

        # Twitter Cards
        for tag in ["twitter:title", "twitter:description", "twitter:image", "twitter:card", "twitter:site"]:
            metadata["html_meta"][tag] = get_meta(name=tag)

        # Language and canonical
        metadata["html_meta"]["lang"] = soup.html.get("lang") if soup.html else None
        canonical = head.find("link", rel="canonical")
        metadata["html_meta"]["canonical"] = canonical["href"] if canonical and canonical.has_attr("href") else None

        # Favicon
        icon = head.find("link", rel=lambda x: x and 'icon' in x.lower())
        metadata["html_meta"]["favicon"] = icon["href"] if icon and icon.has_attr("href") else None

        # Hreflang
        metadata["html_meta"]["hreflangs"] = [
            link["href"] for link in head.find_all("link", rel="alternate")
            if link.has_attr("hreflang") and link.has_attr("href")
        ]

        # Forms
        forms = soup.find_all("form")
        metadata["html_meta"]["forms"] = len(forms)
        metadata["html_meta"]["login_form"] = any("password" in str(form).lower() for form in forms)

        # Media content counts
        metadata["html_meta"]["image_count"] = len(soup.find_all("img"))
        metadata["html_meta"]["video_count"] = len(soup.find_all("video"))
        metadata["html_meta"]["audio_count"] = len(soup.find_all("audio"))
        metadata["html_meta"]["script_count"] = len(soup.find_all("script"))

        # Links
        metadata["html_meta"]["internal_links"] = [
            a["href"] for a in soup.find_all("a", href=True)
            if urlparse(a["href"]).netloc in ["", parsed.netloc]
        ]
        metadata["html_meta"]["external_links"] = [
            a["href"] for a in soup.find_all("a", href=True)
            if urlparse(a["href"]).netloc not in ["", parsed.netloc]
        ]

    except Exception as e:
        metadata["error"] = str(e)

    return metadata



def extract_url_metadata(url: str, timeout: int = 10) -> dict:
    headers = {
        "User-Agent": "Mozilla/5.0 (compatible; MetadataBot/1.0)"
    }
    metadata = {
        "url": url,
        "timestamp": datetime.utcnow().isoformat(),
        "request": {},
        "response": {},
        "headers": {},
        "html_meta": {}
    }

    try:
        response = requests.get(url, headers=headers, timeout=timeout, allow_redirects=True)

        metadata["request"]["final_url"] = response.url
        metadata["request"]["status_code"] = response.status_code
        metadata["request"]["redirects"] = [r.url for r in response.history]
        metadata["response"]["content_type"] = response.headers.get("Content-Type", "")
        metadata["response"]["content_length"] = response.headers.get("Content-Length")
        metadata["response"]["server"] = response.headers.get("Server")
        metadata["response"]["date"] = response.headers.get("Date")
        metadata["headers"] = dict(response.headers)

        # Parse HTML for <head> metadata
        soup = BeautifulSoup(response.text, "html.parser")
        head = soup.head or soup

        def get_meta(name):
            tag = head.find("meta", attrs={"name": name})
            return tag["content"].strip() if tag and tag.has_attr("content") else None

        def get_og(name):
            tag = head.find("meta", attrs={"property": name})
            return tag["content"].strip() if tag and tag.has_attr("content") else None

        metadata["html_meta"]["title"] = (head.title.string.strip()
                                          if head.title and head.title.string else None)
        metadata["html_meta"]["description"] = get_meta("description")
        metadata["html_meta"]["keywords"] = get_meta("keywords")
        metadata["html_meta"]["author"] = get_meta("author")
        metadata["html_meta"]["charset"] = (head.find("meta", charset=True) or {}).get("charset")
        metadata["html_meta"]["viewport"] = get_meta("viewport")
        metadata["html_meta"]["lang"] = (soup.html.get("lang") if soup.html else None)

        # Open Graph / Twitter card
        metadata["html_meta"]["og:title"] = get_og("og:title")
        metadata["html_meta"]["og:description"] = get_og("og:description")
        metadata["html_meta"]["og:image"] = get_og("og:image")
        metadata["html_meta"]["twitter:title"] = get_meta("twitter:title")
        metadata["html_meta"]["twitter:description"] = get_meta("twitter:description")

    except Exception as e:
        metadata["error"] = str(e)

    return metadata


def extract_file_metadata(
    file_path: str
) -> Dict[str, Any]:
    """
    Extract comprehensive metadata from any file type.

    Args:
        file_path (str): Path to target file.

    Returns:
        Dict[str, Any]: Nested metadata structure.
    """
    metadata: Dict[str, Any] = {}
    path = clean_path(file_path)
    if not path:
        return {"error": "File not found"}

    try:
        stats = os.stat(path)
        metadata["size_bytes"] = stats.st_size
        metadata["created"] = datetime.fromtimestamp(stats.st_ctime).isoformat()
        metadata["modified"] = datetime.fromtimestamp(stats.st_mtime).isoformat()
        metadata["mime"] = magic.from_file(path, mime=True)
        
        # Calculate multiple hash types
        with open(path, 'rb') as f:
            content = f.read()
            metadata["hashes"] = {
                "md5": hashlib.md5(content).hexdigest(),
                "sha1": hashlib.sha1(content).hexdigest(),
                "sha256": hashlib.sha256(content).hexdigest()
            }
        
        # Get extended file attributes where supported
        if hasattr(os, 'listxattr'):
            try:
                xattrs = os.listxattr(path)
                if xattrs:
                    metadata["xattrs"] = {}
                    for attr in xattrs:
                        metadata["xattrs"][attr] = os.getxattr(path, attr)
            except (OSError, AttributeError):
                pass
        
        # Get EXIF data if available and relevant
        exif = extract_exif(path)
        if exif:
            metadata["exif"] = exif
            
        # Get file owner and permissions
        try:
            metadata["owner"] = pwd.getpwuid(stats.st_uid).pw_name
        except KeyError:
            metadata["owner"] = str(stats.st_uid)
        metadata["permissions"] = oct(stats.st_mode)[-3:]
                
    except Exception as e:
        metadata["error"] = str(e)
        
    return metadata

def extract_metadata(source: str) -> dict:
    """
    Central metadata extraction router.  Routes metadata extraction based on the source.

    Args:
        source (str): source to extract metadata from

    Returns:
        dict: Dictionary of extracted metadata for the given source.
    """
    if is_url(source):
        metadata = extract_url_metadata(source)
    else:
        metadata = extract_file_metadata(source)

    return metadata

def detect_language(text: str) -> str:
    """
    Detect the language of the extracted text.
    
    Args:
        text (str): Input text
        
    Returns:
        str: Detected language code or 'unknown'
    """
    languages = list_available_languages()
    language_names = {code: name for name, code in languages.items()}
    try:
        import langdetect
        lang = langdetect.detect(text)
        return language_names[lang]
    except:
        logger.warning("Language detection failed or langdetect not installed")
        return "unknown"


def list_available_languages() -> Dict[str, str]:
    """
    Get a dictionary of available languages for translation.

    Returns:
        Dict[str, str]: Dictionary mapping language codes to language names
    """
    try:
        # Get available languages from the translator
        languages = GoogleTranslator().get_supported_languages(as_dict=True)
        return languages
    except Exception as e:
        logger.error(f"Error getting language list: {e}")
        # Return a small subset as fallback
        return {
            "en": "English",
            "es": "Spanish",
            "fr": "French",
            "de": "German",
            "it": "Italian",
            "ja": "Japanese",
            "ko": "Korean",
            "zh-cn": "Chinese (Simplified)",
            "ru": "Russian",
            "ar": "Arabic"
        }

def translate_text(text: str, target_lang: str = "en") -> Optional[str]:
    """
    Translate text to target language.
    
    Args:
        text (str): Input text to translate
        target_lang (str): Target language code (e.g., 'en', 'es', 'fr', 'ja' for Japanese)
        
    Returns:
        Optional[str]: Translated text or None on failure
    """
    try:
        
        # Handle long texts by splitting into chunks (Google has a limit)
        max_chunk_size = 4500  # Google Translate has a limit around 5000 chars
        chunks = []
        
        # Split text into chunks of appropriate size (at sentence boundaries if possible)
        text_remaining = text
        while len(text_remaining) > 0:
            if len(text_remaining) <= max_chunk_size:
                chunks.append(text_remaining)
                break
                
            # Try to find a sentence boundary near the max chunk size
            chunk_end = max_chunk_size
            while chunk_end > 0 and text_remaining[chunk_end] not in ['.', '!', '?', '\n']:
                chunk_end -= 1
                
            # If no good sentence boundary found, just use max size
            if chunk_end == 0:
                chunk_end = max_chunk_size
            else:
                chunk_end += 1  # Include the period or boundary character
                
            chunks.append(text_remaining[:chunk_end])
            text_remaining = text_remaining[chunk_end:]
            
        # Translate each chunk and combine
        translated_chunks = []
        for chunk in chunks:
            translated_chunk = GoogleTranslator(source='auto', target=target_lang).translate(chunk)
            translated_chunks.append(translated_chunk)
            
        return ' '.join(translated_chunks)
    except Exception as e:
        logger.error(f"Translation error: {e}")
        return None


def summarize_text(text: str, sentences: int = 5) -> str:
    """
    Create a simple extractive summary from the text.
    
    Args:
        text (str): Input text to summarize
        sentences (int): Number of sentences to include
        
    Returns:
        str: Summarized text
    """
    try:
       
        # Tokenize and calculate word frequencies
        stop_words = set(stopwords.words('english'))
        sentences_list = sent_tokenize(text)
        
        # If there are fewer sentences than requested, return all
        if len(sentences_list) <= sentences:
            return text
        
        word_frequencies = {}
        for sentence in sentences_list:
            for word in nltk.word_tokenize(sentence):
                word = word.lower()
                if word not in stop_words:
                    if word not in word_frequencies:
                        word_frequencies[word] = 1
                    else:
                        word_frequencies[word] += 1
        
        # Normalize frequencies
        maximum_frequency = max(word_frequencies.values()) if word_frequencies else 1
        for word in word_frequencies:
            word_frequencies[word] = word_frequencies[word] / maximum_frequency
        
        # Score sentences
        sentence_scores = {}
        for i, sentence in enumerate(sentences_list):
            for word in nltk.word_tokenize(sentence.lower()):
                if word in word_frequencies:
                    if i not in sentence_scores:
                        sentence_scores[i] = word_frequencies[word]
                    else:
                        sentence_scores[i] += word_frequencies[word]
        
        # Get top N sentences
        summary_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)[:sentences]
        summary_sentences = [sentences_list[i] for i, _ in sorted(summary_sentences)]
        
        return ' '.join(summary_sentences)
    except Exception as e:
        logger.error(f"Summarization error: {e}")
        return text


def analyze_text(
    text: str,
    advanced: bool = False,
    domain_specific: str = None
) -> Dict[str, Any]:
    """
    Perform comprehensive text analytics with advanced NLP techniques.
    
    Args:
        text (str): Input text
        advanced (bool): Whether to perform computationally intensive advanced analysis
        domain_specific (str): Optional domain for specialized analysis (e.g., "academic", "social_media", "customer_reviews")
        
    Returns:
        Dict: Comprehensive analysis results
    """
    try:
        # Import required libraries
        import numpy as np
        from scipy.spatial import distance
        import networkx as nx
        from sklearn.feature_extraction.text import TfidfVectorizer, CountVectorizer
        from sklearn.decomposition import LatentDirichletAllocation, TruncatedSVD
        from sklearn.cluster import KMeans
       
        # Basic tokenization
        original_text = text
        # Save original case for NER and other cases where case matters
        original_words = nltk.word_tokenize(text)
        original_sentences = nltk.sent_tokenize(text)
        
        # Convert to lowercase for most analysis
        text = text.lower()
        sentences = nltk.sent_tokenize(text)
        words = nltk.word_tokenize(text)
        
        # Filter out punctuation for word-based analysis
        words_no_punct = [word for word in words if word.isalpha()]
        
        # Get paragraphs (text blocks separated by two newlines)
        paragraphs = text.split('\n\n')
        paragraphs = [p.strip() for p in paragraphs if p.strip()]
        
        # Additional paragraph detection for different formats
        if len(paragraphs) <= 1:
            # Try other common paragraph separators
            paragraphs = re.split(r'\n[\t ]*\n', text)
            paragraphs = [p.strip() for p in paragraphs if p.strip()]
            
            # If still only one paragraph, try to detect paragraph by indentation
            if len(paragraphs) <= 1:
                paragraphs = re.split(r'\n[\t ]+', text)
                paragraphs = [p.strip() for p in paragraphs if p.strip()]
        
        # Stopwords
        try:
            stop_words = set(stopwords.words('english'))
        except:
            nltk.download('stopwords')
            stop_words = set(stopwords.words('english'))
        
        # Remove stopwords
        filtered_words = [word for word in words_no_punct if word not in stop_words]
        
        # Stemming and Lemmatization
        stemmer = PorterStemmer()
        lemmatizer = WordNetLemmatizer()
        
        stemmed_words = [stemmer.stem(word) for word in filtered_words]
        lemmatized_words = [lemmatizer.lemmatize(word) for word in filtered_words]
        
        # Word frequencies
        word_freq = Counter(words_no_punct)
        filtered_word_freq = Counter(filtered_words)
        
        # N-grams generation
        bigrams = list(ngrams(words_no_punct, 2))
        trigrams = list(ngrams(words_no_punct, 3))
        fourgrams = list(ngrams(words_no_punct, 4))
        
        bigram_freq = Counter(bigrams)
        trigram_freq = Counter(trigrams)
        fourgram_freq = Counter(fourgrams)
        
        # Part-of-speech tagging
        pos_tags = nltk.pos_tag(original_words)
        pos_counts = Counter([tag for word, tag in pos_tags])
        
        # Count specific parts of speech
        noun_count = sum(1 for _, tag in pos_tags if tag.startswith('NN'))
        verb_count = sum(1 for _, tag in pos_tags if tag.startswith('VB'))
        adj_count = sum(1 for _, tag in pos_tags if tag.startswith('JJ'))
        adv_count = sum(1 for _, tag in pos_tags if tag.startswith('RB'))
        
        # Lexical density (content words / total words)
        content_pos_tags = ['NN', 'VB', 'JJ', 'RB']  # Base forms
        content_words = sum(1 for _, tag in pos_tags if any(tag.startswith(pos) for pos in content_pos_tags))
        lexical_density = content_words / len(words) if words else 0
        
        # Named Entity Recognition
        named_entities = {}
        entity_counts = {}
        
        try:
            ne_chunks = nltk.ne_chunk(pos_tags)
            
            # Process tree to extract named entities
            for chunk in ne_chunks:
                if hasattr(chunk, 'label'):
                    entity_type = chunk.label()
                    entity_text = ' '.join(c[0] for c in chunk.leaves())
                    
                    if entity_type not in named_entities:
                        named_entities[entity_type] = []
                    
                    named_entities[entity_type].append(entity_text)
            
            # Count entities by type
            entity_counts = {entity_type: len(entities) for entity_type, entities in named_entities.items()}
            
        except Exception as ne_error:
            logger.warning(f"NER error: {ne_error}")
            named_entities = {}
            entity_counts = {}
        
        # Basic readability metrics
        char_count = len(text)
        char_count_no_spaces = len(text.replace(" ", ""))
        spaces = char_count - char_count_no_spaces
        
        avg_word_length = sum(len(word) for word in words_no_punct) / len(words_no_punct) if words_no_punct else 0
        avg_sent_length = len(words_no_punct) / len(sentences) if sentences else 0
        avg_para_length = sum(len(p.split()) for p in paragraphs) / len(paragraphs) if paragraphs else 0
        
        # Calculate syllables (approximation)
        def count_syllables(word):
            word = word.lower()
            if len(word) <= 3:
                return 1
            
            # Remove silent e
            if word.endswith('e'):
                word = word[:-1]
                
            # Count vowel groups
            vowels = "aeiouy"
            count = 0
            prev_is_vowel = False
            
            for char in word:
                is_vowel = char in vowels
                if is_vowel and not prev_is_vowel:
                    count += 1
                prev_is_vowel = is_vowel
                
            return max(1, count)  # Return at least 1 syllable
        
        syllable_counts = [count_syllables(word) for word in words_no_punct]
        total_syllables = sum(syllable_counts)
        
        # Readability formulas
        # Flesch Reading Ease
        flesch_reading_ease = 206.835 - (1.015 * avg_sent_length) - (84.6 * (total_syllables / len(words_no_punct))) if words_no_punct else 0
        
        # Flesch-Kincaid Grade Level
        fk_grade = 0.39 * avg_sent_length + 11.8 * (total_syllables / len(words_no_punct)) - 15.59 if words_no_punct else 0
        
        # Gunning Fog Index
        complex_words = sum(1 for word in words_no_punct if count_syllables(word) >= 3)
        complex_word_percentage = complex_words / len(words_no_punct) if words_no_punct else 0
        gunning_fog = 0.4 * (avg_sent_length + 100 * complex_word_percentage) if words_no_punct else 0
        
        # SMOG Index
        if len(sentences) >= 30:
            smog_sentences = sentences[:30]  # Use first 30 sentences
        else:
            smog_sentences = sentences  # Use all available
            
        smog_words = [word for sent in smog_sentences for word in nltk.word_tokenize(sent) if word.isalpha()]
        smog_complex_words = sum(1 for word in smog_words if count_syllables(word) >= 3)
        smog_index = 1.043 * math.sqrt(smog_complex_words * (30 / len(smog_sentences)) if smog_sentences else 0) + 3.1291
        
        # Dale-Chall Readability Formula
        # This would require a list of common words, simplified version:
        dale_chall_diff_words = sum(1 for word in words_no_punct if len(word) >= 7)
        dale_chall_score = 0.1579 * (dale_chall_diff_words / len(words_no_punct) * 100 if words_no_punct else 0) + 0.0496 * avg_sent_length
        
        if dale_chall_diff_words / len(words_no_punct) > 0.05 if words_no_punct else 0:
            dale_chall_score += 3.6365
        
        # Sentiment Analysis
        blob = TextBlob(original_text)
        sentiment = blob.sentiment
        
        # Subjectivity by sentence
        sentence_sentiments = [TextBlob(sent).sentiment for sent in original_sentences]
        sentence_polarities = [sent.polarity for sent in sentence_sentiments]
        sentence_subjectivities = [sent.subjectivity for sent in sentence_sentiments]
        
        # Sentiment variance
        polarity_variance = np.var(sentence_polarities) if sentence_polarities else 0
        subjectivity_variance = np.var(sentence_subjectivities) if sentence_subjectivities else 0
        
        # Sentiment extremes
        most_positive_sentence = original_sentences[np.argmax(sentence_polarities)] if sentence_polarities else ""
        most_negative_sentence = original_sentences[np.argmin(sentence_polarities)] if sentence_polarities else ""
        most_subjective_sentence = original_sentences[np.argmax(sentence_subjectivities)] if sentence_subjectivities else ""
        most_objective_sentence = original_sentences[np.argmin(sentence_subjectivities)] if sentence_subjectivities else ""

        # Averaged categorical sentiment
        positive_threshold = 0.05
        negative_threshold = -0.05
        positive_count = sum(1 for polarity in sentence_polarities if polarity > positive_threshold)
        negative_count = sum(1 for polarity in sentence_polarities if polarity < negative_threshold)
        neutral_count = sum(1 for polarity in sentence_polarities if positive_threshold >= polarity >= negative_threshold)

        # Calculate percentages
        total_sentences = len(sentence_polarities) if sentence_polarities else 1  # Avoid division by zero
        positive_percentage = (positive_count / total_sentences) * 100
        negative_percentage = (negative_count / total_sentences) * 100
        neutral_percentage = (neutral_count / total_sentences) * 100

        # Determine categorical sentiment
        if positive_percentage > 60:
            categorical_sentiment = "very positive"
        elif positive_percentage > 40:
            categorical_sentiment = "positive"
        elif negative_percentage > 60:
            categorical_sentiment = "very negative"
        elif negative_percentage > 40:
            categorical_sentiment = "negative"
        elif neutral_percentage > 60:
            categorical_sentiment = "neutral"
        else:
            categorical_sentiment = "mixed"
        
        # Lexical diversity
        lexical_diversity = len(set(words_no_punct)) / len(words_no_punct) if words_no_punct else 0
        
        # Calculate TF-IDF
        # Without a corpus this is simplified, but we can treat each sentence as a document
        if len(sentences) > 3:  # Only compute if we have enough sentences
            try:
                tfidf_vectorizer = TfidfVectorizer(stop_words='english')
                tfidf_matrix = tfidf_vectorizer.fit_transform(sentences)
                feature_names = tfidf_vectorizer.get_feature_names_out()
                
                # Get top tfidf terms for each sentence
                tfidf_top_terms = []
                for i, sentence in enumerate(sentences):
                    if i < tfidf_matrix.shape[0]:  # Safety check
                        feature_index = tfidf_matrix[i,:].nonzero()[1]
                        tfidf_scores = zip(feature_index, [tfidf_matrix[i, x] for x in feature_index])
                        tfidf_scores = sorted(tfidf_scores, key=lambda x: x[1], reverse=True)
                        tfidf_top_terms.append([(feature_names[i], score) for i, score in tfidf_scores[:5]])
            except Exception as tfidf_error:
                logger.warning(f"TF-IDF error: {tfidf_error}")
                tfidf_top_terms = []
        else:
            tfidf_top_terms = []
        
        # Text summarization - extractive (simplified)
        # Rank sentences by importance (using word frequency as proxy)
        sentence_scores = {}
        for i, sentence in enumerate(sentences):
            sentence_words = nltk.word_tokenize(sentence.lower())
            sentence_words = [word for word in sentence_words if word.isalpha()]
            score = sum(word_freq.get(word, 0) for word in sentence_words)
            sentence_scores[original_sentences[i]] = score
        
        top_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)[:3]
        summary = ' '.join([s[0] for s in top_sentences])
        
        # Text cohesion metrics
        transitional_words = ["however", "therefore", "furthermore", "consequently", "nevertheless", 
                             "thus", "meanwhile", "indeed", "moreover", "whereas", "conversely",
                             "similarly", "in addition", "in contrast", "specifically", "especially",
                             "particularly", "for example", "for instance", "in conclusion", "finally"]
        
        # Count transitional words and their positions
        transition_count = 0
        transition_positions = []
        
        for i, word in enumerate(words):
            if word in transitional_words or any(phrase in ' '.join(words[i:i+4]) for phrase in transitional_words if ' ' in phrase):
                transition_count += 1
                transition_positions.append(i / len(words) if words else 0)  # Normalized position
        
        # Cohesion score - higher means more transitional elements
        cohesion_score = (transition_count / len(words) * 100) if words else 0
        
        # Distribution of transitions (beginning, middle, end)
        if transition_positions:
            transitions_beginning = sum(1 for pos in transition_positions if pos < 0.33)
            transitions_middle = sum(1 for pos in transition_positions if 0.33 <= pos < 0.66)
            transitions_end = sum(1 for pos in transition_positions if pos >= 0.66)
        else:
            transitions_beginning = transitions_middle = transitions_end = 0
            
        # Additional advanced metrics if requested
        advanced_results = {}
        
        if advanced:
            try:
                # Create a document-term matrix
                if len(sentences) >= 5:  # Need enough sentences for meaningful topics
                    # Create Count Vectorizer
                    count_vectorizer = CountVectorizer(stop_words='english', min_df=2)
                    count_matrix = count_vectorizer.fit_transform(sentences)
                    count_feature_names = count_vectorizer.get_feature_names_out()
                    
                    # Train LDA model if we have enough data
                    if count_matrix.shape[0] >= 5 and count_matrix.shape[1] >= 10:
                        n_topics = min(3, count_matrix.shape[0] - 1)  # Choose appropriate number of topics
                        lda_model = LatentDirichletAllocation(n_components=n_topics, random_state=42)
                        lda_model.fit(count_matrix)
                        
                        # Get top words for each topic
                        topics = []
                        for topic_idx, topic in enumerate(lda_model.components_):
                            top_words_idx = topic.argsort()[:-11:-1]  # Top 10 words
                            top_words = [count_feature_names[i] for i in top_words_idx]
                            topics.append(top_words)
                            
                        advanced_results["topics"] = topics
                        
                        # Alternative: Use TruncatedSVD (similar to LSA) for topic extraction
                        svd_model = TruncatedSVD(n_components=n_topics, random_state=42)
                        svd_model.fit(count_matrix)
                        
                        # Get top words for each component (topic)
                        svd_topics = []
                        for topic_idx, topic in enumerate(svd_model.components_):
                            top_words_idx = topic.argsort()[:-11:-1]  # Top 10 words
                            top_words = [count_feature_names[i] for i in top_words_idx]
                            svd_topics.append(top_words)
                            
                        advanced_results["svd_topics"] = svd_topics
                    else:
                        advanced_results["topics"] = ["Insufficient data for topic modeling"]
                        advanced_results["svd_topics"] = ["Insufficient data for topic modeling"]
                else:
                    advanced_results["topics"] = ["Insufficient data for topic modeling"]
                    advanced_results["svd_topics"] = ["Insufficient data for topic modeling"]
                    
                # Clustering sentences instead of document similarity
                if len(sentences) >= 5:
                    # Use TF-IDF vectors for clustering
                    tfidf_vectorizer = TfidfVectorizer(stop_words='english')
                    tfidf_matrix = tfidf_vectorizer.fit_transform(sentences)
                    
                    # Determine number of clusters
                    n_clusters = min(3, len(sentences) - 1)
                    km = KMeans(n_clusters=n_clusters, random_state=42)
                    km.fit(tfidf_matrix)
                    
                    # Get sentence clusters
                    clusters = km.labels_.tolist()
                    
                    # Organize sentences by cluster
                    sentence_clusters = defaultdict(list)
                    for i, cluster in enumerate(clusters):
                        sentence_clusters[cluster].append(original_sentences[i])
                        
                    advanced_results["sentence_clusters"] = dict(sentence_clusters)
                    
                    # Get top terms per cluster
                    cluster_terms = {}
                    order_centroids = km.cluster_centers_.argsort()[:, ::-1]
                    terms = tfidf_vectorizer.get_feature_names_out()
                    
                    for i in range(n_clusters):
                        cluster_top_terms = [terms[ind] for ind in order_centroids[i, :10]]
                        cluster_terms[i] = cluster_top_terms
                        
                    advanced_results["cluster_terms"] = cluster_terms
                else:
                    advanced_results["sentence_clusters"] = {"note": "Insufficient data for clustering"}
                    advanced_results["cluster_terms"] = {"note": "Insufficient data for clustering"}
                
            except Exception as topic_error:
                logger.warning(f"Topic modeling error: {topic_error}")
                advanced_results["topics"] = ["Error in topic modeling"]
                advanced_results["svd_topics"] = ["Error in topic modeling"]
            
            # Text network analysis
            try:
                # Create word co-occurrence network
                G = nx.Graph()
                
                # Add nodes (words)
                for word in set(filtered_words):
                    G.add_node(word)
                
                # Add edges (co-occurrences within sentences)
                for sentence in sentences:
                    sent_words = [word for word in nltk.word_tokenize(sentence.lower()) 
                                 if word.isalpha() and word not in stop_words]
                    
                    # Add edges between all pairs of words in the sentence
                    for i, word1 in enumerate(sent_words):
                        for word2 in sent_words[i+1:]:
                            if G.has_edge(word1, word2):
                                G[word1][word2]['weight'] += 1
                            else:
                                G.add_edge(word1, word2, weight=1)
                
                # Calculate network metrics if we have enough data
                if G.number_of_nodes() > 2:
                    # Degree centrality
                    degree_centrality = nx.degree_centrality(G)
                    top_degree_centrality = sorted(degree_centrality.items(), key=lambda x: x[1], reverse=True)[:10]
                    
                    # Betweenness centrality for central connector words
                    if G.number_of_nodes() < 1000:  # Skip for very large networks
                        betweenness_centrality = nx.betweenness_centrality(G, k=min(G.number_of_nodes(), 100))
                        top_betweenness = sorted(betweenness_centrality.items(), key=lambda x: x[1], reverse=True)[:10]
                    else:
                        top_betweenness = [("Network too large", 0)]
                    
                    # Extract clusters/communities (simplified)
                    components = list(nx.connected_components(G))
                    largest_component = max(components, key=len)
                    
                    advanced_results["network_analysis"] = {
                        "central_terms": [word for word, score in top_degree_centrality],
                        "connector_terms": [word for word, score in top_betweenness],
                        "clusters_count": len(components),
                        "largest_cluster_size": len(largest_component)
                    }
                else:
                    advanced_results["network_analysis"] = {"note": "Insufficient data for network analysis"}
            except Exception as network_error:
                logger.warning(f"Network analysis error: {network_error}")
                advanced_results["network_analysis"] = {"error": str(network_error)}
                
            # Syntactic complexity
            try:
                # Parse subtrees (approximation)
                syntactic_complexity = {}
                
                # Count depth of clauses (approximation using POS patterns)
                clause_markers = [',', 'that', 'which', 'who', 'whom', 'whose', 'where', 'when', 'why', 'how']
                subordinating_conjunctions = ['after', 'although', 'as', 'because', 'before', 'if', 'since', 'though', 'unless', 'until', 'when', 'where', 'while']
                
                clause_complexity = []
                
                for sentence in original_sentences:
                    tokens = nltk.word_tokenize(sentence.lower())
                    clause_markers_count = sum(1 for token in tokens if token in clause_markers)
                    subordinating_count = sum(1 for token in tokens if token in subordinating_conjunctions)
                    
                    # Approximate clause depth
                    clause_depth = 1 + clause_markers_count + subordinating_count
                    clause_complexity.append(clause_depth)
                
                syntactic_complexity["avg_clause_depth"] = sum(clause_complexity) / len(clause_complexity) if clause_complexity else 0
                syntactic_complexity["max_clause_depth"] = max(clause_complexity) if clause_complexity else 0
                
                # Approximation of phrase types
                sentence_pos = [nltk.pos_tag(nltk.word_tokenize(sentence)) for sentence in original_sentences]
                
                # Count noun phrases (approximated by adjective-noun sequences)
                noun_phrases = []
                for sentence_tags in sentence_pos:
                    for i in range(len(sentence_tags) - 1):
                        if sentence_tags[i][1].startswith('JJ') and sentence_tags[i+1][1].startswith('NN'):
                            noun_phrases.append(f"{sentence_tags[i][0]} {sentence_tags[i+1][0]}")
                
                # Count verb phrases (approximated by adverb-verb sequences)
                verb_phrases = []
                for sentence_tags in sentence_pos:
                    for i in range(len(sentence_tags) - 1):
                        if sentence_tags[i][1].startswith('RB') and sentence_tags[i+1][1].startswith('VB'):
                            verb_phrases.append(f"{sentence_tags[i][0]} {sentence_tags[i+1][0]}")
                
                syntactic_complexity["estimated_noun_phrases"] = len(noun_phrases)
                syntactic_complexity["estimated_verb_phrases"] = len(verb_phrases)
                syntactic_complexity["noun_verb_phrase_ratio"] = len(noun_phrases) / len(verb_phrases) if verb_phrases else 0
                
                advanced_results["syntactic_complexity"] = syntactic_complexity
            except Exception as syntax_error:
                logger.warning(f"Syntactic analysis error: {syntax_error}")
                advanced_results["syntactic_complexity"] = {"error": str(syntax_error)}
                
            # Domain-specific analysis
            if domain_specific:
                domain_analysis = {}
                
                if domain_specific == "academic":
                    # Academic writing analysis
                    academic_terms = ["hypothesis", "theory", "analysis", "data", "method", "research", 
                                    "study", "evidence", "results", "conclusion", "findings", "literature",
                                    "significant", "therefore", "thus", "however", "moreover"]
                    hedge_words = ["may", "might", "could", "appears", "seems", "suggests", "indicates",
                                  "possibly", "perhaps", "likely", "unlikely", "generally", "usually"]
                    
                    academic_term_count = sum(word_freq.get(term, 0) for term in academic_terms)
                    hedge_word_count = sum(word_freq.get(term, 0) for term in hedge_words)
                    
                    domain_analysis["academic_term_density"] = academic_term_count / len(words) if words else 0
                    domain_analysis["hedging_density"] = hedge_word_count / len(words) if words else 0
                    domain_analysis["citation_count"] = original_text.count("et al") + re.findall(r"\(\d{4}\)", original_text).__len__()
                    
                elif domain_specific == "social_media":
                    # Social media analysis
                    hashtags = re.findall(r"#\w+", original_text)
                    mentions = re.findall(r"@\w+", original_text)
                    urls = re.findall(r"https?://\S+", original_text)
                    emojis = re.findall(r"[\U0001F600-\U0001F64F\U0001F300-\U0001F5FF\U0001F680-\U0001F6FF\U0001F700-\U0001F77F\U0001F780-\U0001F7FF\U0001F800-\U0001F8FF\U0001F900-\U0001F9FF\U0001FA00-\U0001FA6F\U0001FA70-\U0001FAFF\U00002702-\U000027B0\U000024C2-\U0001F251]+", original_text)
                    
                    slang_terms = ["lol", "omg", "wtf", "idk", "tbh", "imo", "fwiw", "ymmv", "tl;dr", "ftw"]
                    slang_count = sum(1 for word in words if word.lower() in slang_terms)
                    
                    domain_analysis["hashtag_count"] = len(hashtags)
                    domain_analysis["mention_count"] = len(mentions)
                    domain_analysis["url_count"] = len(urls)
                    domain_analysis["emoji_count"] = len(emojis)
                    domain_analysis["slang_terms"] = slang_count
                    domain_analysis["engagement_markers"] = len(hashtags) + len(mentions) + len(emojis) + slang_count
                    
                elif domain_specific == "customer_reviews":
                    # Customer review analysis
                    product_terms = ["product", "quality", "price", "value", "recommend", "purchase",
                                   "buy", "bought", "worth", "money", "shipping", "delivery", "package",
                                   "arrived", "customer", "service", "return", "warranty", "replacement"]
                    
                    rating_terms = ["star", "stars", "rating", "rate", "perfect", "excellent", "good", 
                                  "average", "poor", "terrible", "worst", "best"]
                    
                    feature_terms = ["feature", "features", "works", "worked", "functionality", "design",
                                   "size", "weight", "color", "material", "battery", "screen", "interface"]
                    
                    product_term_count = sum(word_freq.get(term, 0) for term in product_terms)
                    rating_term_count = sum(word_freq.get(term, 0) for term in rating_terms)
                    feature_term_count = sum(word_freq.get(term, 0) for term in feature_terms)
                    
                    # Find potential ratings (e.g. "5 star", "3.5 out of 5")
                    rating_patterns = re.findall(r"(\d+\.?\d*)\s*(star|stars|out of \d+)", original_text.lower())
                    
                    domain_analysis["product_term_density"] = product_term_count / len(words) if words else 0
                    domain_analysis["rating_term_density"] = rating_term_count / len(words) if words else 0
                    domain_analysis["feature_term_density"] = feature_term_count / len(words) if words else 0
                    domain_analysis["potential_ratings"] = rating_patterns
                    domain_analysis["recommendation_language"] = "recommend" in text.lower() or "would buy" in text.lower()
                
                advanced_results["domain_analysis"] = domain_analysis
        
        # Collect all results
        analysis_results = {
            # Basic counts
            "basic_stats": {
                "word_count": len(words),
                "unique_word_count": len(set(words_no_punct)),
                "sentence_count": len(sentences),
                "paragraph_count": len(paragraphs),
                "character_count": char_count,
                "character_count_no_spaces": char_count_no_spaces,
                "avg_word_length": avg_word_length,
                "avg_sentence_length": avg_sent_length,
                "avg_paragraph_length": avg_para_length,
                "spaces": spaces,
                "punctuation_count": len(original_words) - len(words_no_punct),
            },
                        # Part of speech
            "part_of_speech": {
                "distribution": dict(pos_counts),
                "noun_count": noun_count,
                "verb_count": verb_count,
                "adjective_count": adj_count,
                "adverb_count": adv_count,
                "noun_to_verb_ratio": noun_count / verb_count if verb_count else 0,
                "lexical_density": lexical_density,
            },

            # Named entities
            "named_entities": {
                "counts": entity_counts,
                "entities": named_entities,
                "total_entities": sum(entity_counts.values())
            },

            # Readability metrics
            "readability": {
                "flesch_reading_ease": flesch_reading_ease,
                "flesch_kincaid_grade": fk_grade,
                "gunning_fog_index": gunning_fog,
                "smog_index": smog_index,
                "dale_chall_score": dale_chall_score,
                "syllable_count": total_syllables,
                "avg_syllables_per_word": total_syllables / len(words_no_punct) if words_no_punct else 0,
                "complex_word_percentage": complex_word_percentage,
                "lexical_diversity": lexical_diversity,
            },

            # Frequency analysis
            "frequency_analysis": {
                "most_common_words": word_freq.most_common(20),
                "most_common_meaningful_words": filtered_word_freq.most_common(20),
                "most_common_bigrams": bigram_freq.most_common(10),
                "most_common_trigrams": trigram_freq.most_common(5),
                "most_common_fourgrams": fourgram_freq.most_common(3),
                "hapax_legomena": [word for word, count in word_freq.items() if count == 1],  # Words occurring only once
                "hapax_percentage": sum(1 for _, count in word_freq.items() if count == 1) / len(word_freq) if word_freq else 0,
            },

            # Sentiment analysis
            "sentiment": {
                "overall_polarity": sentiment.polarity,  # -1 to 1 (negative to positive)
                "overall_subjectivity": sentiment.subjectivity,  # 0 to 1 (objective to subjective)
                "polarity_variance": polarity_variance,
                "subjectivity_variance": subjectivity_variance,
                "most_positive_sentence": most_positive_sentence,
                "most_negative_sentence": most_negative_sentence,
                "most_subjective_sentence": most_subjective_sentence,
                "most_objective_sentence": most_objective_sentence,
                "sentiment_shifts": sum(1 for i in range(1, len(sentence_polarities))
                                      if (sentence_polarities[i-1] > 0 and sentence_polarities[i] < 0) or
                                         (sentence_polarities[i-1] < 0 and sentence_polarities[i] > 0)),
                "sentiment_progression": "positive_trend" if sum(1 for i in range(1, len(sentence_polarities))
                                                           if sentence_polarities[i] > sentence_polarities[i-1]) > len(sentence_polarities) / 2
                                     else "negative_trend" if sum(1 for i in range(1, len(sentence_polarities))
                                                           if sentence_polarities[i] < sentence_polarities[i-1]) > len(sentence_polarities) / 2
                                     else "neutral_trend",
                "categorical_sentiment": {
                    "label": categorical_sentiment,
                    "positive_percentage": positive_percentage,
                    "negative_percentage": negative_percentage,
                    "neutral_percentage": neutral_percentage,
                    "positive_sentence_count": positive_count,
                    "negative_sentence_count": negative_count,
                    "neutral_sentence_count": neutral_count
                }
            },

            # Preprocessing results
            "preprocessing": {
                "filtered_words_count": len(filtered_words),
                "stopwords_removed": len(words_no_punct) - len(filtered_words),
                "stemmed_words_sample": stemmed_words[:10] if stemmed_words else [],
                "lemmatized_words_sample": lemmatized_words[:10] if lemmatized_words else [],
            },

            # Text summarization
            "summarization": {
                "extractive_summary": summary,
                "key_sentences": [s[0] for s in top_sentences],
                "tfidf_top_terms": tfidf_top_terms
            },

            # Text cohesion metrics
            "cohesion": {
                "transitional_word_count": transition_count,
                "cohesion_score": cohesion_score,
                "transitions_beginning": transitions_beginning,
                "transitions_middle": transitions_middle,
                "transitions_end": transitions_end,
                "connector_distribution": "front_loaded" if transitions_beginning > transitions_middle and transitions_beginning > transitions_end
                                      else "end_loaded" if transitions_end > transitions_beginning and transitions_end > transitions_middle
                                      else "evenly_distributed"
            }
        }

        # Add advanced results if they were computed
        if advanced and advanced_results:
            for key, value in advanced_results.items():
                analysis_results[key] = value

        analysis_results["language"] = detect_language(text)

        # Topic modeling (simple keyword-based approach)
        topic_keywords = filtered_word_freq.most_common(10)
        analysis_results["topic_analysis"] = {
            "possible_topics": topic_keywords
        }

        # Contextual analysis (identifying context patterns)
        contextual_analysis = {}

        # Temporal references
        temporal_markers = ["today", "yesterday", "tomorrow", "now", "then", "before", "after",
                           "while", "during", "soon", "later", "earlier", "recently", "ago"]
        temporal_references = sum(word_freq.get(marker, 0) for marker in temporal_markers)

        # Spatial references
        spatial_markers = ["here", "there", "above", "below", "behind", "in front", "nearby",
                          "inside", "outside", "around", "between", "among", "everywhere"]
        spatial_references = sum(word_freq.get(marker, 0) for marker in spatial_markers)

        # Personal references
        first_person = sum(word_freq.get(marker, 0) for marker in ["i", "me", "my", "mine", "we", "us", "our", "ours"])
        second_person = sum(word_freq.get(marker, 0) for marker in ["you", "your", "yours"])
        third_person = sum(word_freq.get(marker, 0) for marker in ["he", "him", "his", "she", "her", "hers", "they", "them", "their", "theirs"])

        contextual_analysis["temporal_references"] = temporal_references
        contextual_analysis["spatial_references"] = spatial_references
        contextual_analysis["first_person_references"] = first_person
        contextual_analysis["second_person_references"] = second_person
        contextual_analysis["third_person_references"] = third_person
        contextual_analysis["narration_perspective"] = "first_person" if first_person > second_person and first_person > third_person else \
                                              "second_person" if second_person > first_person and second_person > third_person else \
                                              "third_person"

        analysis_results["contextual_analysis"] = contextual_analysis

        # Detect writing style (tentative classification)
        style_markers = {}

        # Formality markers
        formal_markers = ["therefore", "thus", "consequently", "furthermore", "moreover", "hence",
                        "accordingly", "subsequently", "previously", "regarding", "concerning"]
        informal_markers = ["anyway", "basically", "actually", "kinda", "like", "so", "pretty",
                          "totally", "really", "hopefully", "maybe", "ok", "okay", "stuff"]

        style_markers["formal_marker_count"] = sum(word_freq.get(marker, 0) for marker in formal_markers)
        style_markers["informal_marker_count"] = sum(word_freq.get(marker, 0) for marker in informal_markers)
        style_markers["contraction_count"] = len(re.findall(r"\b\w+'[ts]|\b\w+n't\b|\b\w+'ll\b|\b\w+'re\b|\b\w+'ve\b", original_text))
        style_markers["exclamation_count"] = original_text.count("!")
        style_markers["question_count"] = original_text.count("?")
        style_markers["parenthetical_count"] = len(re.findall(r"\([^)]*\)", original_text))
        style_markers["semicolon_count"] = original_text.count(";")

        # Tentative style classification
        formality_score = (style_markers["formal_marker_count"] + style_markers["semicolon_count"] * 2 +
                          style_markers["parenthetical_count"]) - (style_markers["informal_marker_count"] +
                          style_markers["contraction_count"] + style_markers["exclamation_count"] * 2)

        if formality_score > 5:
            style_markers["style_classification"] = "formal_academic"
        elif formality_score > 0:
            style_markers["style_classification"] = "formal"
        elif formality_score > -5:
            style_markers["style_classification"] = "neutral"
        elif formality_score > -10:
            style_markers["style_classification"] = "informal"
        else:
            style_markers["style_classification"] = "very_informal"

        analysis_results["style_analysis"] = style_markers

        # Detect potential rhetoric patterns
        rhetoric_patterns = {}

        # Repetition patterns
        repeated_bigrams = [bg for bg, count in bigram_freq.items() if count > 2]
        repeated_trigrams = [tg for tg, count in trigram_freq.items() if count > 2]

        # Question patterns
        rhetorical_questions = sum(1 for sentence in original_sentences if
                                 sentence.endswith("?") and any(word in sentence.lower() for word in
                                                              ["why", "who", "what", "how", "when", "where"]))

        # Comparison patterns (similes)
        similes = len(re.findall(r"\b(like|as) a\b|\b(like|as) the\b", original_text.lower()))

        # Alliteration (simplified detection)
        alliterations = 0
        for i in range(len(original_words) - 2):
            if (len(original_words[i]) > 0 and len(original_words[i+1]) > 0 and len(original_words[i+2]) > 0 and
                original_words[i][0].lower() == original_words[i+1][0].lower() == original_words[i+2][0].lower()):
                alliterations += 1

        rhetoric_patterns["repeated_phrases"] = repeated_bigrams + repeated_trigrams
        rhetoric_patterns["rhetorical_questions"] = rhetorical_questions
        rhetoric_patterns["similes"] = similes
        rhetoric_patterns["alliterations"] = alliterations

        analysis_results["rhetoric_patterns"] = rhetoric_patterns

        # Potential bias indicators
        bias_indicators = {}

        # Extreme language
        extreme_markers = ["always", "never", "all", "none", "every", "only", "impossible",
                          "absolutely", "undoubtedly", "certainly", "definitely", "completely",
                          "total", "totally", "utterly", "best", "worst", "perfect"]
        extreme_language = sum(word_freq.get(marker, 0) for marker in extreme_markers)

        # Loaded language
        emotionally_loaded = ["amazing", "terrible", "awesome", "horrible", "wonderful", "dreadful",
                             "excellent", "awful", "extraordinary", "appalling", "incredible", "disgusting"]
        loaded_language = sum(word_freq.get(marker, 0) for marker in emotionally_loaded)

        bias_indicators["extreme_language_count"] = extreme_language
        bias_indicators["loaded_language_count"] = loaded_language
        bias_indicators["extreme_language_ratio"] = extreme_language / len(words) if words else 0
        bias_indicators["loaded_language_ratio"] = loaded_language / len(words) if words else 0

        # Simple bias classification
        if bias_indicators["extreme_language_ratio"] > 0.05 or bias_indicators["loaded_language_ratio"] > 0.05:
            bias_indicators["bias_classification"] = "potentially_biased"
        else:
            bias_indicators["bias_classification"] = "relatively_neutral"

        analysis_results["bias_indicators"] = bias_indicators

        if advanced and len(filtered_words) >= 20:
            try:
                similarity_analysis = {}

                # Create co-occurrence matrix (simplified word embedding alternative)
                vectorizer = CountVectorizer(ngram_range=(1, 1), token_pattern=r'\b\w+\b', min_df=2)
                X = vectorizer.fit_transform(sentences)
                features = vectorizer.get_feature_names_out()

                # Convert to array for easier manipulation
                X_array = X.toarray()

                # Compute pairwise distance between terms
                # Use cosine similarity between term vectors
                term_similarity = {}

                for i, term1 in enumerate(features):
                    if i < len(X_array[0]):  # Safety check
                        term_vec1 = X_array[:, i]
                        for j, term2 in enumerate(features):
                            if j < len(X_array[0]) and i != j:  # Skip self-comparison
                                term_vec2 = X_array[:, j]
                                # Compute cosine similarity
                                similarity = 1 - distance.cosine(term_vec1, term_vec2)

                                if similarity > 0.5:  # Only keep high similarity pairs
                                    if term1 not in term_similarity:
                                        term_similarity[term1] = []
                                    term_similarity[term1].append((term2, similarity))

                # Sort and keep top similar terms
                for term in term_similarity:
                    term_similarity[term] = sorted(term_similarity[term], key=lambda x: x[1], reverse=True)[:5]

                # Get top terms with most connections
                top_connected_terms = sorted(term_similarity.items(), key=lambda x: len(x[1]), reverse=True)[:10]

                similarity_analysis["term_similarity"] = {term: similar for term, similar in top_connected_terms}

                advanced_results["similarity_analysis"] = similarity_analysis
            except Exception as sim_error:
                logger.warning(f"Similarity analysis error: {sim_error}")
                advanced_results["similarity_analysis"] = {"error": str(sim_error)}

        # Emotion detection (beyond just sentiment)
        try:
            emotion_analysis = {}

            # Basic emotion lexicons
            emotions = {
                "joy": ["happy", "joy", "delight", "glad", "pleased", "excited", "thrilled", "elated"],
                "sadness": ["sad", "unhappy", "sorrow", "depressed", "miserable", "downcast", "gloomy"],
                "anger": ["angry", "mad", "furious", "outraged", "annoyed", "irritated", "livid"],
                "fear": ["afraid", "fear", "scared", "terrified", "worried", "anxious", "nervous"],
                "surprise": ["surprised", "amazed", "astonished", "shocked", "stunned", "startled"],
                "disgust": ["disgusted", "revolted", "repulsed", "sickened", "appalled"]
            }

            # Count emotion words
            emotion_counts = {}
            for emotion, emotion_words in emotions.items():
                emotion_counts[emotion] = sum(word_freq.get(word, 0) for word in emotion_words)

            # Calculate dominant emotion
            dominant_emotion = max(emotion_counts.items(), key=lambda x: x[1]) if emotion_counts else ("neutral", 0)

            # Calculate emotion intensity (percent of all emotion words that belong to dominant emotion)
            total_emotion_words = sum(emotion_counts.values())
            dominant_intensity = (dominant_emotion[1] / total_emotion_words) if total_emotion_words > 0 else 0

            emotion_analysis["emotion_counts"] = emotion_counts
            emotion_analysis["dominant_emotion"] = dominant_emotion[0]
            emotion_analysis["dominant_intensity"] = dominant_intensity
            emotion_analysis["emotional_diversity"] = len([e for e, c in emotion_counts.items() if c > 0])

            analysis_results["emotion_analysis"] = emotion_analysis
        except Exception as emo_error:
            logger.warning(f"Emotion analysis error: {emo_error}")
            analysis_results["emotion_analysis"] = {"error": str(emo_error)}

        # Get metadata about the analysis
        metadata = {
            "timestamp": datetime.now().isoformat(),
            "analysis_version": "2.0",
            "text_length_category": "short" if len(words) < 100 else "medium" if len(words) < 500 else "long",
            "advanced_analysis_performed": advanced,
            "domain_specific_analysis": domain_specific
        }

        analysis_results["metadata"] = metadata

        return analysis_results

    except Exception as e:
        logger.error(f"Text analysis error: {e}")
        import traceback
        return {
            "error": str(e),
            "traceback": traceback.format_exc()
        }

def text_from_epub(epub_path: str) -> Optional[str]:
    """
    Extract text from EPUB ebooks.
    
    Args:
        epub_path (str): Path to the EPUB file
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        from ebooklib import epub
        import html2text
        
        book = epub.read_epub(epub_path)
        h = html2text.HTML2Text()
        h.ignore_links = False
        
        content = []
        
        # Get book metadata
        metadata = []
        if book.get_metadata('DC', 'title'):
            metadata.append(f"Title: {book.get_metadata('DC', 'title')[0][0]}")
        if book.get_metadata('DC', 'creator'):
            metadata.append(f"Author: {book.get_metadata('DC', 'creator')[0][0]}")
        if book.get_metadata('DC', 'description'):
            metadata.append(f"Description: {book.get_metadata('DC', 'description')[0][0]}")
        
        if metadata:
            content.append("\n".join(metadata))
            content.append("---")
        
        # Get book content
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                content.append(h.handle(item.get_content().decode('utf-8')))
        
        return normalize_text("\n".join(content))
    except ImportError:
        logger.error("ebooklib and/or html2text not installed")
        return "ebooklib and/or html2text packages are required for EPUB processing"
    except Exception as e:
        logger.error(f"Error processing EPUB: {e}")
        return None


def text_from_pptx(pptx_path: str) -> Optional[str]:
    """
    Extract text from PowerPoint presentations.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        text = ["--- PowerPoint Presentation ---"]
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {i}:"]
            
            # Get slide title if it exists
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # Extract text from all shapes
            shape_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    shape_text.append(shape.text)
            
            if shape_text:
                slide_text.append("\n".join(shape_text))
            
            text.append("\n".join(slide_text))
        
        return normalize_text("\n\n".join(text))
    except ImportError:
        logger.error("python-pptx not installed")
        return "python-pptx package is required for PowerPoint processing"
    except Exception as e:
        logger.error(f"Error processing PowerPoint: {e}")
        return None


def text_from_odt(odt_path: str) -> Optional[str]:
    """
    Extract text from OpenDocument Text files.
    
    Args:
        odt_path (str): Path to the ODT file
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        from odf import text, teletype
        from odf.opendocument import load
        
        textdoc = load(odt_path)
        
        # Extract metadata
        meta = []
        meta_elem = textdoc.meta
        if meta_elem:
            for prop in meta_elem.childNodes:
                if hasattr(prop, 'tagName') and hasattr(prop, 'childNodes') and prop.childNodes:
                    meta.append(f"{prop.tagName}: {teletype.extractText(prop)}")
        
        # Extract content
        allparas = textdoc.getElementsByType(text.P)
        content = "\n".join(teletype.extractText(p) for p in allparas)
        
        # Combine metadata and content
        if meta:
            final_text = "\n".join(meta) + "\n---\n" + content
        else:
            final_text = content
        
        return normalize_text(final_text)
    except ImportError:
        logger.error("odfpy not installed")
        return "odfpy package is required for ODT processing"
    except Exception as e:
        logger.error(f"Error processing ODT: {e}")
        return None


# Memory optimization functions

def extract_text_chunked(file_path: str, chunk_size: int = 10) -> Optional[str]:
    """
    Extract text from very large files in chunks to manage memory usage.
    
    Args:
        file_path (str): Path to the large file
        chunk_size (int): Number of pages/sections to process at once
        
    Returns:
        Optional[str]: Combined text
    """
    path = clean_path(file_path)
    if not path:
        return None
        
    mime_type = magic.from_file(path, mime=True)
    
    if mime_type == "application/pdf":
        return extract_pdf_chunked(path, chunk_size)
    elif mime_type.startswith("text/"):
        return extract_text_file_chunked(path, chunk_size)
    else:
        return extract_text(path)  # Fall back to regular extraction


def extract_text_file_chunked(file_path: str, chunk_size: int = 1024*1024) -> str:
    """
    Extract text from large text files in chunks.
    
    Args:
        file_path (str): Path to the text file
        chunk_size (int): Size of each chunk in bytes
        
    Returns:
        str: Combined text
    """
    try:
        result = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            while True:
                chunk = f.read(chunk_size)
                if not chunk:
                    break
                result.append(chunk)
        
        return normalize_text("".join(result))
    except Exception as e:
        logger.error(f"Error processing text file in chunks: {e}")
        return ""


def batch_extract(file_paths: List[str], max_workers: Optional[int] = None) -> Dict[str, str]:
    """
    Process multiple files concurrently.
    
    Args:
        file_paths (List[str]): List of file paths to process
        max_workers (int, optional): Maximum number of worker threads
        
    Returns:
        Dict[str, str]: Dictionary mapping file paths to extracted text
    """
    results = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {executor.submit(extract_text, file_path): file_path for file_path in file_paths}
        for future in concurrent.futures.as_completed(future_to_file):
            file_path = future_to_file[future]
            try:
                results[file_path] = future.result() or f"No text extracted from {file_path}"
            except Exception as exc:
                results[file_path] = f"ERROR: {exc}"
    return results


def process_extracted_text(text: str, args) -> Union[str, Dict[str, Any]]:
    """
    Process extracted text according to command-line arguments.
    
    Args:
        text (str): Extracted text
        args: Command-line arguments
        
    Returns:
        Union[str, Dict[str, Any]]: Processed text or analysis results
    """
    # Process text according to arguments
    if args.summarize:
        logger.info(f"Summarizing text (sentences: {args.sentences})")
        text = summarize_text(text, args.sentences)
    
    if args.analyze:
        logger.info("Analyzing text")
        analysis = analyze_text(text)
        return analysis
    
    if args.translate and args.translate != "list":
        logger.info(f"Translating text to {args.translate}")
        translated = translate_text(text, args.translate)
        if translated:
            text = translated
        else:
            logger.error("Translation failed")
    
    return text


def main() -> None:
    """
    Enhanced CLI entry point for text extraction, metadata, and analytics.
    Supports both file/URL inputs and stdin piped data.
    """
    import os
    import argparse
    import glob
    import tempfile
    import sys
    from io import BytesIO
    
    class TranslateAction(argparse.Action):
        def __call__(self, parser, namespace, values, option_string=None):
            # If no value is provided, set to 'list' to trigger language listing
            if values is None:
                setattr(namespace, self.dest, 'list')
            else:
                setattr(namespace, self.dest, values)
    
    parser = argparse.ArgumentParser(
        description="Extract and analyze text from any file, URL, directory, wildcard pattern, or stdin"
    )
    parser.add_argument(
        "source",
        nargs="*",  # Make source completely optional
        help="Path(s) to file(s), URL, directory, or wildcard pattern (omit to read from stdin)"
    )
    parser.add_argument(
        "--metadata",
        action="store_true",
        help="Extract metadata instead of text"
    )
    parser.add_argument(
        "--summarize",
        action="store_true",
        help="Summarize the extracted text"
    )
    parser.add_argument(
        "--sentences",
        type=int,
        default=5,
        help="Number of sentences in summary (default: 5)"
    )
    parser.add_argument(
        "--analyze",
        action="store_true",
        help="Perform text analysis"
    )
    parser.add_argument(
        "--translate",
        action=TranslateAction,
        nargs="?",  # Allow --translate with no argument to list languages
        metavar="LANG",
        help="Translate text to specified language code (e.g., 'es'), or list available languages if no code provided"
    )
    parser.add_argument(
        "--output",
        type=str,
        help="Output file path (default: stdout)"
    )
    parser.add_argument(
        "--password",
        type=str,
        help="Password for protected documents"
    )
    parser.add_argument(
        "--scrape",
        action="store_true",
        help="Scrape multiple pages from a website (for URLs only)"
    )
    parser.add_argument(
        "--max-pages",
        type=int,
        default=5,
        help="Maximum pages to scrape when using --scrape (default: 5)"
    )
    parser.add_argument(
        "--verbose", "-v",
        action="count",
        default=0,
        help="Increase verbosity (can be used multiple times)"
    )
    parser.add_argument(
        "--screenshot",
        action="store_true",
        help="Capture and extract text from screen"
    )
    parser.add_argument(
        "--chunked",
        action="store_true",
        help="Process large files in chunks to reduce memory usage"
    )
    parser.add_argument(
        "--no-js",
        action="store_true",
        help="Disable JavaScript rendering for web pages"
    )
    parser.add_argument(
        "--list-languages",
        action="store_true",
        help="List available translation languages"
    )
    parser.add_argument(
        "--stdin-format",
        type=str,
        help="Specify input format when reading from stdin (pdf, docx, txt, etc.)"
    )
    
    args = parser.parse_args()
    
    # Configure logging level based on verbosity
    if args.verbose == 0:
        logger.setLevel(logging.WARNING)
    elif args.verbose == 1:
        logger.setLevel(logging.INFO)
    else:
        logger.setLevel(logging.DEBUG)
    
    # Handle translation language listing (in two ways for better usability)
    if args.translate == "list" or args.list_languages:
        try:
            logger.info("Fetching available languages...")
            languages = list_available_languages()
            
            print("\n[Available Translation Languages]")
            # Format the output as a table with columns
            max_code_len = max(len(code) for code in languages.keys())
            max_name_len = max(len(name) for name in languages.values())
            format_str = f"  {{:<{max_code_len}}}  {{:<{max_name_len}}}"
            
            print(format_str.format("Code", "Language"))
            print(format_str.format("-" * max_code_len, "-" * max_name_len))
            
            for code, name in sorted(languages.items(), key=lambda x: x[1]):  # Sort by language name
                print(format_str.format(code, name))
            
            sys.exit(0)
        except Exception as e:
            logger.error(f"Error listing languages: {e}")
            sys.exit(1)
    
    # Function to handle output
    def output_result(result):
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                if isinstance(result, dict):
                    json.dump(result, f, indent=2, ensure_ascii=False)
                else:
                    f.write(str(result))
            logger.info(f"Output written to {args.output}")
        else:
            if isinstance(result, dict):
                print(json.dumps(result, indent=2, ensure_ascii=False))
            else:
                print(result)
    
    # Handle screenshot mode
    if args.screenshot:
        logger.info("Capturing screenshot...")
        text = text_from_screenshot()
        if not text:
            logger.error("No text extracted from screenshot.")
            sys.exit(1)
        
        # Process the text according to arguments
        result = process_extracted_text(text, args)
        output_result(result)
        sys.exit(0)
    
    # Check if we should read from stdin (no source arguments provided and not in screenshot mode)
    if not args.source and not args.screenshot and not sys.stdin.isatty():
        logger.info("Reading data from stdin...")
        
        # Create a temporary file to store stdin data
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            try:
                # Read binary data from stdin
                stdin_data = sys.stdin.buffer.read()
                
                if not stdin_data:
                    logger.error("No data received from stdin.")
                    sys.exit(1)
                
                # Write the data to the temporary file
                temp_file.write(stdin_data)
                temp_path = temp_file.name
                
                # Close the file to ensure all data is written
                temp_file.close()
                
                # Determine the file type if not explicitly provided
                mime_type = None
                if args.stdin_format:
                    # User specified the format
                    format_map = {
                        'pdf': 'application/pdf',
                        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                        'doc': 'application/msword',
                        'txt': 'text/plain',
                        'html': 'text/html',
                        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                        'xls': 'application/vnd.ms-excel',
                        'json': 'application/json',
                        'xml': 'application/xml',
                    }
                    mime_type = format_map.get(args.stdin_format.lower())
                
                # If format wasn't specified or wasn't in our map, try to detect it
                if not mime_type:
                    import magic
                    mime_type = magic.from_file(temp_path, mime=True)
                    logger.info(f"Detected MIME type: {mime_type}")
                
                # Extract text based on the detected or specified format
                if args.metadata:
                    logger.info("Extracting metadata from stdin data...")
                    metadata = extract_metadata(temp_path)
                    output_result(metadata)
                else:
                    if args.chunked:
                        logger.info("Extracting text from stdin data in chunks...")
                        text = extract_text_chunked(temp_path)
                    elif args.password:
                        logger.info("Extracting text from password-protected stdin data...")
                        text = extract_text_with_password(temp_path, args.password)
                    else:
                        logger.info("Extracting text from stdin data...")
                        text = extract_text(temp_path)
                    
                    if not text:
                        logger.error("No text extracted from stdin data.")
                        sys.exit(1)
                    
                    result = process_extracted_text(text, args)
                    output_result(result)
                
            except Exception as e:
                logger.error(f"Error processing stdin data: {e}")
                sys.exit(1)
            finally:
                # Clean up the temporary file
                try:
                    import os
                    os.unlink(temp_path)
                except Exception as e:
                    logger.error(f"Failed to delete temp file: {e}")
        
        sys.exit(0)
    
    # Process sources (if provided)
    sources = args.source
    
    # If no sources and stdin is a tty (terminal), show help
    if not sources and sys.stdin.isatty() and not args.screenshot:
        parser.print_help()
        sys.exit(0)
    
    # Expand any directory sources
    files = []
    for source in sources:
        if os.path.isdir(source):
            # It's a directory, add all files in it
            logger.info(f"Processing directory: {source}")
            dir_files = [os.path.join(source, f) for f in os.listdir(source) 
                        if os.path.isfile(os.path.join(source, f))]
            files.extend(dir_files)
        elif any(c in source for c in ["*", "?", "[", "]"]) and not os.path.exists(source):
            # It's a wildcard pattern that wasn't expanded by the shell
            logger.info(f"Processing pattern: {source}")
            expanded = glob.glob(source, recursive=True)
            if expanded:
                files.extend(expanded)
            else:
                logger.warning(f"No files found matching pattern: {source}")
        elif os.path.isfile(source):
            # It's a regular file
            files.append(source)
        elif is_url(source):
            # It's a URL - we'll handle it separately
            if len(sources) == 1:
                # Only process URL if it's the only source
                if args.scrape:
                    logger.info(f"Scraping website: {source} (max {args.max_pages} pages)")
                    results = scrape_website(source, max_pages=args.max_pages)
                    combined_text = "\n\n".join([f"=== {url} ===\n{text}" for url, text in results.items()])
                    result = process_extracted_text(combined_text, args)
                    output_result(result)
                else:
                    logger.info(f"Extracting text from URL: {source}")
                    text = text_from_url(source, render_js=not args.no_js)
                    if not text:
                        logger.error("No text extracted.")
                        sys.exit(1)
                    result = process_extracted_text(text, args)
                    output_result(result)
                return
            else:
                logger.warning(f"Skipping URL {source} when processing multiple sources")
        else:
            logger.warning(f"Source not found: {source}")
    
    if not files:
        logger.error("No valid files found to process.")
        sys.exit(1)
    
    logger.info(f"Processing {len(files)} files")
    
    # Process metadata if requested
    if args.metadata:
        results = {}
        with concurrent.futures.ThreadPoolExecutor() as executor:
            future_to_file = {executor.submit(extract_metadata, file_path): file_path for file_path in files}
            for future in concurrent.futures.as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    results[os.path.basename(file_path)] = future.result()
                except Exception as exc:
                    results[os.path.basename(file_path)] = {"error": str(exc)}
        output_result(results)
        return
    
    # If only one file and not metadata, process it directly
    if len(files) == 1:
        file_path = files[0]
        try:
            if args.chunked:
                logger.info(f"Extracting text from {file_path} in chunks")
                text = extract_text_chunked(file_path)
            elif args.password:
                logger.info(f"Extracting text from password-protected {file_path}")
                text = extract_text_with_password(file_path, args.password)
            else:
                logger.info(f"Extracting text from {file_path}")
                text = extract_text(file_path)
                
            if not text:
                logger.error("No text extracted.")
                sys.exit(1)
                
            result = process_extracted_text(text, args)
            output_result(result)
            return
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            sys.exit(1)
    
    # Process multiple files
    results = batch_extract(files)
    
    # Process each extracted text according to arguments
    processed_results = {}
    for file_path, text in results.items():
        if isinstance(text, str) and not text.startswith("ERROR:"):
            processed_results[os.path.basename(file_path)] = process_extracted_text(text, args)
        else:
            processed_results[os.path.basename(file_path)] = text
    
    output_result(processed_results)


# Execute as script
if __name__ == "__main__":
    main()    
