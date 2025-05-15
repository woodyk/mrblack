#!/usr/bin/env python3
# -*- coding: utf-8 -*-
#
# File: textextract.py
# Author: Wadih Khairallah
# Created: 2024-12-01 12:12:08
# Modified: 2025-05-15 16:18:26
# Enhanced with additional features

import os
import re
import sys
import json
import math
import socket
import platform
import subprocess
import hashlib
import string
import tempfile
import random
import unicodedata
import concurrent.futures
import logging
from logging import Logger

from uuid import uuid4
from datetime import datetime
from io import StringIO
from typing import (
    Optional,
    Dict,
    Any,
    List,
    Union,
    Tuple,
    Set,
    Callable,
)

# Web Specific
import requests
from requests_html import HTMLSession
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup

# Document Specific
from readability import Document as RDocument
import magic
import pytesseract
import pandas as pd
import speech_recognition as sr
import pymupdf
from docx import Document
from mss import mss
from pydub import AudioSegment

# Image Specific
from PIL import Image

# NLP and Text Analytics
from collections import Counter
import nltk
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords', quiet=True)

# Logging with Rich
from rich.console import Console
from rich.logging import RichHandler

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format="%(message)s",
    datefmt="[%X]",
    handlers=[RichHandler(rich_tracebacks=True)]
)
logger = logging.getLogger("textextract")

console = Console()

# Use rich for pretty printing by default
print = console.print
log = logger.info

# User agent strings for web requests
USER_AGENTS = [
    # Desktop
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 13_3) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Safari/605.1.15",
    "Mozilla/5.0 (X11; Linux x86_64) Gecko/20100101 Firefox/115.0",
    # Mobile
    "Mozilla/5.0 (iPhone; CPU iPhone OS 16_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Mobile/15E148 Safari/604.1",
    "Mozilla/5.0 (Linux; Android 13; Pixel 7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.6367.91 Mobile Safari/537.36"
]


def clean_path(
    path: str
) -> Optional[str]:
    """
    Normalize and validate a filesystem path.

    Args:
        path (str): Input file or directory path.

    Returns:
        Optional[str]: Absolute path if valid; None otherwise.
    """
    p = os.path.expanduser(path)
    p = os.path.abspath(p)
    if os.path.isfile(p) or os.path.isdir(p):
        return p
    return None


def normalize(
    text: str
) -> str:
    """
    Replace multiple consecutive newlines, carriage returns, and spaces
    with a single space. Ensures compact, single-line output.

    Args:
        text (str): Raw input text.

    Returns:
        str: Normalized single-line text.
    """
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r' +', ' ', text)
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r'(?m)(^ \n)+', '\n', text)
    text = re.sub(r'\t+', '\t', text)
    text = re.sub(r'\r+', '\n', text)
    text = re.sub(r"^ ", "", text, flags=re.MULTILINE)
    return text 


def is_url(s: str) -> bool:
    """
    Check if a string is a valid HTTP/HTTPS URL.
    
    Args:
        s (str): String to check
        
    Returns:
        bool: True if valid URL, False otherwise
    """
    parsed = urlparse(s)
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def text_from_screenshot() -> str:
    """
    Capture a full-screen screenshot, perform OCR, and clean up temp file.

    Returns:
        str: Normalized OCR-extracted text from the screenshot.
    """
    tmp_filename = f"screenshot_{uuid4().hex}.png"
    tmp_path = os.path.join(tempfile.gettempdir(), tmp_filename)

    try:
        with mss() as sct:
            monitor = {"top": 0, "left": 0, "width": 0, "height": 0}
            for mon in sct.monitors:
                monitor["left"] = min(mon["left"], monitor["left"])
                monitor["top"] = min(mon["top"], monitor["top"])
                monitor["width"] = max(mon["width"] + mon["left"] - monitor["left"], monitor["width"])
                monitor["height"] = max(mon["height"] + mon["top"] - monitor["top"], monitor["height"])
            screenshot = sct.grab(monitor)

        img = Image.frombytes("RGB", screenshot.size, screenshot.bgra, "raw", "BGRX")
        img_gray = img.convert("L")
        img_gray.save(tmp_path)

        content = text_from_image(tmp_path)
        return normalize(content)
    finally:
        if os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except Exception as e:
                logger.error(f"Failed to delete temp screenshot: {e}")


def extract_exif(
    file_path: str
) -> Optional[Dict[str, Any]]:
    """
    Extract EXIF metadata from a file using exiftool.

    Args:
        file_path (str): Path to the target file.

    Returns:
        Optional[Dict[str, Any]]: Parsed EXIF data, or None on failure.
    """
    exif_data: Optional[Dict[str, Any]] = None
    try:
        result = subprocess.run(
            ['exiftool', '-j', file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        if result.returncode == 0:
            exif_data = json.loads(result.stdout.decode())[0]
    except Exception as e:
        logger.error(f"Exiftool failed: {e}")
    return exif_data


def extract_document_structure(html: str) -> Dict[str, List[Dict[str, Any]]]:
    """
    Extract document structure (headings, lists, sections) from HTML.
    
    Args:
        html (str): HTML content
        
    Returns:
        Dict: Document structure information
    """
    soup = BeautifulSoup(html, "html.parser")
    
    structure = {
        "headings": [],
        "lists": [],
        "sections": []
    }
    
    # Extract headings
    for level in range(1, 7):
        for heading in soup.find_all(f'h{level}'):
            structure["headings"].append({
                "level": level,
                "text": heading.get_text(strip=True),
                "id": heading.get("id", "")
            })
    
    # Extract lists
    for list_type in ['ul', 'ol']:
        for lst in soup.find_all(list_type):
            items = [li.get_text(strip=True) for li in lst.find_all('li')]
            structure["lists"].append({
                "type": "unordered" if list_type == "ul" else "ordered",
                "items": items
            })
    
    # Extract sections
    for section in soup.find_all('section'):
        section_info = {
            "id": section.get("id", ""),
            "class": section.get("class", []),
            "text": section.get_text(strip=True)
        }
        structure["sections"].append(section_info)
    
    return structure


def text_from_html(html: str) -> str:
    """
    Extract readable text from raw HTML content.

    Args:
        html (str): HTML source as a string.

    Returns:
        str: Cleaned and normalized visible text.
    """
    # Check if the input is a file path or HTML content
    if os.path.isfile(html):
        with open(html, 'r', encoding='utf-8', errors='ignore') as f:
            html = f.read()
    
    soup = BeautifulSoup(html, "html.parser")

    # Remove non-visible or structural elements
    for tag in soup([
        "script", "style",
        "noscript", "iframe",
        "meta", "link",
        "header", "footer",
        "form", "nav",
        "aside"
    ]):
        tag.decompose()

    text = soup.get_text(separator=" ")

    return normalize(text)


def text_from_url(
    url: str,
    render_js: bool = True
) -> Optional[str]:
    """
    Fetch and extract all visible text from a web page, including JS-rendered content.

    Args:
        url (str): Target webpage URL.
        render_js (bool): Whether to render JavaScript content.

    Returns:
        Optional[str]: Cleaned full-page text, or None on failure.
    """
    headers = {
        "User-Agent": random.choice(USER_AGENTS),
        "Accept-Language": "en-US,en;q=0.9",
        "Referer": url,
        "DNT": "1",
        "Upgrade-Insecure-Requests": "1"
    }

    # Try with requests-html first (with JS rendering)
    if render_js:
        try:
            session = HTMLSession()
            try:
                r = session.get(url, headers=headers, timeout=20)
                
                # Set shorter timeout for rendering to avoid hanging
                try:
                    r.html.render(timeout=10, sleep=1, keep_page=True)
                except Exception as e:
                    logger.warning(f"JS rendering failed, falling back to static HTML: {e}")
                
                html = r.html.html
                session.close()
                content = text_from_html(html)
                return content
            except Exception as e:
                logger.error(f"[Error with HTMLSession] {url} - {e}")
                session.close()
                # Fall through to regular requests
            finally:
                session.close()
        except Exception as e:
            logger.error(f"[Error creating HTMLSession] {e}")
            # Fall through to regular requests
    
    # Fall back to regular requests (no JS rendering)
    try:
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        html = response.text
        content = text_from_html(html)
        return content
    except Exception as e:
        logger.error(f"[Error with requests] {url} - {e}")
        return None


def scrape_website(url: str, max_pages: int = 1, stay_on_domain: bool = True) -> Dict[str, str]:
    """
    Scrape multiple pages of a website.
    
    Args:
        url (str): Starting URL
        max_pages (int): Maximum pages to scrape
        stay_on_domain (bool): Whether to stay on the same domain
        
    Returns:
        Dict[str, str]: Dictionary mapping URLs to extracted text
    """
    results = {}
    visited = set()
    to_visit = [url]
    base_domain = urlparse(url).netloc
    
    while to_visit and len(visited) < max_pages:
        current_url = to_visit.pop(0)
        if current_url in visited:
            continue
            
        # Extract text from current page
        text = text_from_url(current_url)
        if text:
            results[current_url] = text
            
        visited.add(current_url)
        
        # Find links on the page
        session = HTMLSession()
        try:
            r = session.get(current_url)
            r.html.render(timeout=20, sleep=1)
            
            links = r.html.absolute_links
            for link in links:
                link_domain = urlparse(link).netloc
                if link not in visited and link not in to_visit:
                    # Check if we should follow this link
                    if stay_on_domain and link_domain != base_domain:
                        continue
                    to_visit.append(link)
        except Exception as e:
            logger.error(f"Error scraping {current_url}: {e}")
        finally:
            session.close()
    
    return results


def extract_text(
    file_path: str
) -> Optional[str]:
    """
    Extract text content from a local file or URL.

    Supports web pages, text, JSON, XML, CSV, Excel, PDF, DOCX, images, audio.

    Args:
        file_path (str): Path to the input file or URL.

    Returns:
        Optional[str]: Extracted text, or None if unsupported or error.
    """
    if is_url(file_path):
        return text_from_url(file_path)

    TEXT_MIME_TYPES = {
        "application/json", "application/xml", "application/x-yaml",
        "application/x-toml", "application/x-csv", "application/x-markdown",
    }

    path = clean_path(file_path)
    if not path:
        logger.error(f"No such file: {file_path}")
        return None

    mime_type = magic.from_file(path, mime=True)
    try:
        if mime_type.startswith("text/html"):
            content = text_from_html(path)
            return content

        if mime_type.startswith("text/") or mime_type in TEXT_MIME_TYPES:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return normalize(content)

        elif mime_type in [
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ]:
            content = text_from_excel(path)
            return content

        elif mime_type == "application/pdf":
            content = text_from_pdf(path)
            return content

        elif mime_type == \
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            content = text_from_docx(path)
            return content

        elif mime_type == "application/msword":
            content = text_from_doc(path)
            return content

        elif mime_type.startswith("image/"):
            content = text_from_image(path)
            return content

        elif mime_type.startswith("audio/"):
            content = text_from_audio(path)
            return content

        elif mime_type == "application/epub+zip":
            content = text_from_epub(path)
            return content

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            content = text_from_pptx(path)
            return content

        elif mime_type == "application/vnd.oasis.opendocument.text":
            content = text_from_odt(path)
            return content

        else:
            content = text_from_any(path)
            return content
    except Exception as e:
        logger.error(f"Error reading {path}: {e}")
        return None


def extract_text_with_password(file_path: str, password: str) -> Optional[str]:
    """
    Extract text from password-protected files.
    
    Args:
        file_path (str): Path to the file
        password (str): Password to unlock the file
        
    Returns:
        Optional[str]: Extracted text
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.pdf':
        return text_from_pdf_protected(file_path, password)
    elif file_ext in ['.docx', '.xlsx', '.pptx']:
        return text_from_office_protected(file_path, password)
    else:
        logger.warning(f"Password protection not supported for {file_ext} files")
        return None


def text_from_pdf_protected(pdf_path: str, password: str) -> Optional[str]:
    """
    Extract text from password-protected PDF.
    
    Args:
        pdf_path (str): Path to the PDF file
        password (str): Password to unlock the PDF
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        doc = pymupdf.open(pdf_path, password=password)
        text = ""
        for i in range(len(doc)):
            page = doc.load_page(i)
            text += page.get_text()
        doc.close()
        return normalize(text)
    except Exception as e:
        logger.error(f"Error in PDF with password: {str(e)}")
        return None


def text_from_office_protected(office_path: str, password: str) -> Optional[str]:
    """
    Extract text from password-protected Office files.
    
    Args:
        office_path (str): Path to the Office file
        password (str): Password to unlock the file
        
    Returns:
        Optional[str]: Extracted text
    """
    # This is a placeholder. Actual implementation would depend on the specific office format
    # and available libraries for password-protected office files.
    logger.warning("Password-protected Office files extraction not fully implemented")
    return None


def text_from_audio(
    audio_file: str
) -> Optional[str]:
    """
    Transcribe audio to text using Google Speech Recognition.

    Args:
        audio_file (str): Path to the input audio file.

    Returns:
        Optional[str]: Transcription, or None on failure.
    """
    def convert_to_wav(file_path: str) -> str:
        _, ext = os.path.splitext(file_path)
        ext = ext.lstrip('.')
        audio = AudioSegment.from_file(file_path, format=ext)
        tmp_filename = f"audio_{uuid4().hex}.wav"
        wav_path = os.path.join(tempfile.gettempdir(), tmp_filename)
        audio.export(wav_path, format='wav')
        return wav_path

    recognizer = sr.Recognizer()
    temp_wav_path = None
    cleanup_needed = False

    try:
        _, ext = os.path.splitext(audio_file)
        if ext.lower() not in ['.wav', '.wave']:
            temp_wav_path = convert_to_wav(audio_file)
            cleanup_needed = True
        else:
            temp_wav_path = clean_path(audio_file)

        if not temp_wav_path:
            logger.error("Invalid audio path.")
            return None

        with sr.AudioFile(temp_wav_path) as source:
            audio = recognizer.record(source)
        return recognizer.recognize_google(audio)

    except sr.UnknownValueError:
        logger.error("Could not understand audio.")
    except sr.RequestError as e:
        logger.error(f"Speech recognition error: {e}")
    except Exception as e:
        logger.error(f"Failed to process audio: {e}")
    finally:
        if cleanup_needed and temp_wav_path and os.path.exists(temp_wav_path):
            try:
                os.remove(temp_wav_path)
            except Exception as e:
                logger.error(f"Failed to delete temp WAV file {temp_wav_path}: {e}")

    return None


def downloadImage(
    url: str
) -> Optional[str]:
    """
    Download an image from a URL to /tmp/ and return its path.

    Args:
        url (str): Remote image URL.

    Returns:
        Optional[str]: Local file path, or None on failure.
    """
    if is_url(url):
        try:
            resp = requests.head(url)
            content_type = resp.headers.get('Content-Type', '')
            if content_type.startswith('image/'):
                filename = os.path.basename(urlparse(url).path)
                if not filename:
                    filename = f"image_{uuid4().hex}.jpg"  # Default name if none found
                save_path = os.path.join(tempfile.gettempdir(), filename)
                resp = requests.get(url, stream=True)
                resp.raise_for_status()
                with open(save_path, "wb") as f:
                    for chunk in resp.iter_content(chunk_size=8192):
                        f.write(chunk)
                return clean_path(save_path)
        except Exception as e:
            logger.error(f"Failed to download image from {url}: {e}")
    logger.error(f"Unable to pull image from {url}")
    return None


def is_image(
    file_path_or_url: str
) -> bool:
    """
    Determine if the given path/URL points to an image.

    Args:
        file_path_or_url (str): Local path or URL.

    Returns:
        bool: True if MIME type starts with 'image/'.
    """
    try:
        if is_url(file_path_or_url):
            resp = requests.head(file_path_or_url)
            content_type = resp.headers.get('Content-Type', '')
            return content_type.startswith('image/')
        else:
            mime = magic.from_file(file_path_or_url, mime=True)
            return mime.startswith("image/")
    except Exception:
        return False


def text_from_pdf(
    pdf_path: str
) -> Optional[str]:
    """
    Extract text and OCR results from a PDF using PyMuPDF.

    Args:
        pdf_path (str): Path to PDF file.

    Returns:
        Optional[str]: Combined normalized text and image OCR results.
    """
    plain_text = ""
    temp_image_paths: List[str] = []

    try:
        doc = pymupdf.open(pdf_path)
        for k, v in doc.metadata.items():
            plain_text += f"{k}: {v}\n"

        for i in range(len(doc)):
            page = doc.load_page(i)
            plain_text += f"\n--- Page {i + 1} ---\n"
            text = page.get_text()
            plain_text += text or "[No text]\n"

            for img_index, img in enumerate(page.get_images(full=True), start=1):
                xref = img[0]
                base = doc.extract_image(xref)
                img_bytes = base["image"]

                img_filename = f"pdf_page{i+1}_img{img_index}_{uuid4().hex}.png"
                img_path = os.path.join(tempfile.gettempdir(), img_filename)
                temp_image_paths.append(img_path)

                with open(img_path, "wb") as f:
                    f.write(img_bytes)

                ocr = text_from_image(img_path) or ""
                plain_text += f"\n[Image {img_index} OCR]\n{ocr}\n"

        # Extract tables from PDF
        """
        try:
            tables = extract_tables_from_pdf(pdf_path)
            if tables:
                plain_text += "\n--- Tables ---\n"
                for i, table in enumerate(tables, 1):
                    plain_text += f"\n[Table {i}]\n"
                    if isinstance(table, dict) and "data" in table:
                        for row in table["data"]:
                            plain_text += str(row) + "\n"
                    else:
                        plain_text += str(table) + "\n"
        except Exception as e:
            logger.warning(f"Could not extract tables from PDF: {e}")
        """

        return normalize(plain_text)
    except Exception as e:
        logger.error(f"Error processing PDF: {e}")
        return None
    finally:
        for path in temp_image_paths:
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception as e:
                    logger.error(f"Failed to delete temp image {path}: {e}")
        if 'doc' in locals():
            doc.close()


def extract_tables_from_pdf(pdf_path: str) -> List[Dict[str, Any]]:
    """
    Extract tables from PDF using tabula-py.
    
    Args:
        pdf_path (str): Path to the PDF file
        
    Returns:
        List[Dict[str, Any]]: List of extracted tables
    """
    try:
        import tabula
        
        tables = []
        # Read all tables from the PDF
        extracted_tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
        
        for i, table in enumerate(extracted_tables):
            if not table.empty:
                tables.append({
                    "table_number": i+1,
                    "rows": len(table),
                    "columns": len(table.columns),
                    "data": table.to_dict(orient='records')
                })
        
        return tables
    except ImportError:
        logger.warning("tabula-py not installed, skipping table extraction")
        return []
    except Exception as e:
        logger.error(f"Error extracting tables from PDF: {e}")
        return []


def extract_pdf_chunked(pdf_path: str, chunk_size: int = 10) -> Optional[str]:
    """
    Extract text from large PDFs in chunks to manage memory usage.
    
    Args:
        pdf_path (str): Path to the PDF file
        chunk_size (int): Number of pages to process at once
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        doc = pymupdf.open(pdf_path)
        total_pages = len(doc)
        text_chunks = []
        
        # Add metadata
        metadata = ""
        for k, v in doc.metadata.items():
            metadata += f"{k}: {v}\n"
        text_chunks.append(metadata)
        
        for i in range(0, total_pages, chunk_size):
            chunk_text = ""
            end_page = min(i + chunk_size, total_pages)
            
            for page_num in range(i, end_page):
                page = doc.load_page(page_num)
                chunk_text += f"\n--- Page {page_num + 1} ---\n"
                chunk_text += page.get_text() or "[No text]\n"
                
            text_chunks.append(chunk_text)
            
        doc.close()
        return normalize("".join(text_chunks))
    except Exception as e:
        logger.error(f"Error processing PDF in chunks: {e}")
        return None


def text_from_doc(
    filepath: str,
    min_length: int = 4
) -> str:
    """
    Extract readable strings and metadata from binary Word (.doc) files.

    Args:
        filepath (str): Path to .doc file.
        min_length (int): Minimum string length to extract.

    Returns:
        str: Metadata and text content.
    """
    def extract_printable_strings(
        data: bytes
    ) -> List[str]:
        pattern = re.compile(
            b'[' + re.escape(bytes(string.printable, 'ascii')) +
            b']{%d,}' % min_length
        )
        found = pattern.findall(data)
        return list(dict.fromkeys(m.decode(errors='ignore').strip()
                                   for m in found))

    def clean_strings(
        strs: List[str]
    ) -> List[str]:
        cleaned: List[str] = []
        skip = ["HYPERLINK", "OLE2", "Normal.dotm"]
        for line in strs:
            if any(line.startswith(pref) for pref in skip):
                continue
            cleaned.append(re.sub(r'\s+', ' ', line).strip())
        return cleaned

    with open(filepath, 'rb') as f:
        data = f.read()
    strings = extract_printable_strings(data)
    strings = clean_strings(strings)
    content = "\n".join(strings)
    return normalize(content)


def text_from_docx(
    file_path: str
) -> Optional[str]:
    """
    Extract text, tables, and OCR from embedded images in a DOCX file.

    Args:
        file_path (str): Path to the .docx file.

    Returns:
        Optional[str]: Normalized full text content.
    """
    path = clean_path(file_path)
    if not path:
        return None

    temp_image_paths: List[str] = []
    plain_text = ""

    try:
        doc = Document(path)

        for p in doc.paragraphs:
            if p.text.strip():
                plain_text += p.text.strip() + "\n"

        for tbl in doc.tables:
            plain_text += "\n[Table]\n"
            for row in tbl.rows:
                row_text = "\t".join(c.text.strip() for c in row.cells)
                plain_text += row_text + "\n"

        for rel_id, rel in doc.part.rels.items():
            if "image" in rel.target_ref:
                blob = rel.target_part.blob

                img_filename = f"docx_img_{rel_id}_{uuid4().hex}.png"
                img_path = os.path.join(tempfile.gettempdir(), img_filename)
                temp_image_paths.append(img_path)

                with open(img_path, "wb") as img_file:
                    img_file.write(blob)

                ocr = text_from_image(img_path) or ""
                plain_text += f"\n[Image OCR]\n{ocr}\n"

        return normalize(plain_text)

    except Exception as e:
        logger.error(f"Error processing DOCX: {e}")
        return None
    finally:
        for path in temp_image_paths:
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception as e:
                    logger.error(f"Failed to delete temp DOCX image {path}: {e}")


def text_from_excel(
    file_path: str
) -> str:
    """
    Convert an Excel workbook to CSV text.

    Args:
        file_path (str): Path to the Excel file.

    Returns:
        str: CSV-formatted string.
    """
    path = clean_path(file_path)
    if not path:
        return ""
    try:
        # Get all sheets
        result = ""
        excel_file = pd.ExcelFile(path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet_name)
            out = StringIO()
            df.to_csv(out, index=False)
            result += f"\n--- Sheet: {sheet_name} ---\n"
            result += out.getvalue()
            result += "\n"
        return result
    except Exception as e:
        logger.error(f"Failed Excel -> CSV: {e}")
        return ""


def text_from_image(
    file_path: str
) -> Optional[str]:
    """
    Perform OCR on an image file.

    Args:
        file_path (str): Path to the image.

    Returns:
        Optional[str]: Extracted text, or None on error.
    """
    path = clean_path(file_path)
    if not path:
        return None
    try:
        with Image.open(path) as img:
            # Improve OCR with preprocessing
            # 1. Convert to grayscale if it's not already
            if img.mode != 'L':
                img = img.convert('L')
                
            # 2. Optional: Apply some contrast enhancement
            # (Disabled by default, enable if needed for specific cases)
            # from PIL import ImageEnhance
            # enhancer = ImageEnhance.Contrast(img)
            # img = enhancer.enhance(1.5)  # Increase contrast
                
            # Perform OCR with custom configuration
            custom_config = r'--oem 3 --psm 6'  # Default OCR Engine Mode and Page Segmentation Mode
            txt = pytesseract.image_to_string(img, config=custom_config).strip()
            return normalize(txt) or ""
    except Exception as e:
        logger.error(f"Failed image OCR: {e}")
        return None


def text_from_any(
    file_path: str
) -> Optional[str]:
    """
    Handle unknown file types by reporting stats and metadata.

    Args:
        file_path (str): Path to the file.

    Returns:
        Optional[str]: Plain-text report, or None on error.
    """
    path = clean_path(file_path)
    if not path:
        return None
    try:
        stats = os.stat(path)
        info = {
            "path": path,
            "size": stats.st_size,
            "created": datetime.fromtimestamp(stats.st_ctime).isoformat(),
            "modified": datetime.fromtimestamp(stats.st_mtime).isoformat(),
        }
        
        # Try to extract EXIF if available
        exif = extract_exif(path)
        if exif:
            info["exif"] = exif
            
        # Get file hash
        md5_hash = hashlib.md5(open(path,'rb').read()).hexdigest()
        info["md5"] = md5_hash
        
        content = "\n".join(f"{k}: {v}" for k, v in info.items() if k != "exif")
        
        # Add formatted EXIF data if available
        if exif:
            content += "\n\nEXIF Data:\n"
            for k, v in exif.items():
                if isinstance(v, dict):
                    content += f"\n{k}:\n"
                    for sub_k, sub_v in v.items():
                        content += f"  {sub_k}: {sub_v}\n"
                else:
                    content += f"{k}: {v}\n"
                    
        return normalize(content)
    except Exception as e:
        logger.error(f"Error on other file: {e}")
        return None


def extract_metadata(
    file_path: str
) -> Dict[str, Any]:
    """
    Extract comprehensive metadata from any file type.

    Args:
        file_path (str): Path to target file.

    Returns:
        Dict[str, Any]: Nested metadata structure.
    """
    path = clean_path(file_path)
    if not path:
        return {"error": "File not found"}
    
    meta: Dict[str, Any] = {}
    try:
        stats = os.stat(path)
        meta["size_bytes"] = stats.st_size
        meta["created"] = datetime.fromtimestamp(stats.st_ctime).isoformat()
        meta["modified"] = datetime.fromtimestamp(stats.st_mtime).isoformat()
        meta["mime"] = magic.from_file(path, mime=True)
        
        # Calculate multiple hash types
        with open(path, 'rb') as f:
            content = f.read()
            meta["hashes"] = {
                "md5": hashlib.md5(content).hexdigest(),
                "sha1": hashlib.sha1(content).hexdigest(),
                "sha256": hashlib.sha256(content).hexdigest()
            }
        
        # Get extended file attributes where supported
        if hasattr(os, 'listxattr'):
            try:
                xattrs = os.listxattr(path)
                if xattrs:
                    meta["xattrs"] = {}
                    for attr in xattrs:
                        meta["xattrs"][attr] = os.getxattr(path, attr)
            except (OSError, AttributeError):
                pass
        
        # Get EXIF data if available and relevant
        exif = extract_exif(path)
        if exif:
            meta["exif"] = exif
            
        # Get file owner and permissions
        import pwd
        try:
            meta["owner"] = pwd.getpwuid(stats.st_uid).pw_name
        except KeyError:
            meta["owner"] = str(stats.st_uid)
        meta["permissions"] = oct(stats.st_mode)[-3:]
            
    except Exception as e:
        meta["error"] = str(e)
        
    return meta


# New functions for text analytics and processing

def detect_language(text: str) -> str:
    """
    Detect the language of the extracted text.
    
    Args:
        text (str): Input text
        
    Returns:
        str: Detected language code or 'unknown'
    """
    try:
        import langdetect
        return langdetect.detect(text)
    except:
        logger.warning("Language detection failed or langdetect not installed")
        return "unknown"

def list_available_languages() -> Dict[str, str]:
    """
    Get a dictionary of available languages for translation.

    Returns:
        Dict[str, str]: Dictionary mapping language codes to language names
    """
    try:
        from deep_translator import GoogleTranslator
        # Get available languages from the translator
        languages = GoogleTranslator().get_supported_languages(as_dict=True)
        return languages
    except Exception as e:
        logger.error(f"Error getting language list: {e}")
        # Return a small subset as fallback
        return {
            "en": "English",
            "es": "Spanish",
            "fr": "French",
            "de": "German",
            "it": "Italian",
            "ja": "Japanese",
            "ko": "Korean",
            "zh-cn": "Chinese (Simplified)",
            "ru": "Russian",
            "ar": "Arabic"
        }

def translate_text(text: str, target_lang: str = "en") -> Optional[str]:
    """
    Translate text to target language.
    
    Args:
        text (str): Input text to translate
        target_lang (str): Target language code (e.g., 'en', 'es', 'fr', 'ja' for Japanese)
        
    Returns:
        Optional[str]: Translated text or None on failure
    """
    try:
        # Use a more stable translation library
        # Note: googletrans 4.0.0-rc1 uses async methods which need to be awaited
        # Let's use the deep-translator library instead which is more stable
        from deep_translator import GoogleTranslator
        
        # Handle long texts by splitting into chunks (Google has a limit)
        max_chunk_size = 4500  # Google Translate has a limit around 5000 chars
        chunks = []
        
        # Split text into chunks of appropriate size (at sentence boundaries if possible)
        text_remaining = text
        while len(text_remaining) > 0:
            if len(text_remaining) <= max_chunk_size:
                chunks.append(text_remaining)
                break
                
            # Try to find a sentence boundary near the max chunk size
            chunk_end = max_chunk_size
            while chunk_end > 0 and text_remaining[chunk_end] not in ['.', '!', '?', '\n']:
                chunk_end -= 1
                
            # If no good sentence boundary found, just use max size
            if chunk_end == 0:
                chunk_end = max_chunk_size
            else:
                chunk_end += 1  # Include the period or boundary character
                
            chunks.append(text_remaining[:chunk_end])
            text_remaining = text_remaining[chunk_end:]
            
        # Translate each chunk and combine
        translated_chunks = []
        for chunk in chunks:
            translated_chunk = GoogleTranslator(source='auto', target=target_lang).translate(chunk)
            translated_chunks.append(translated_chunk)
            
        return ' '.join(translated_chunks)
    except Exception as e:
        logger.error(f"Translation error: {e}")
        return None


def summarize_text(text: str, sentences: int = 5) -> str:
    """
    Create a simple extractive summary from the text.
    
    Args:
        text (str): Input text to summarize
        sentences (int): Number of sentences to include
        
    Returns:
        str: Summarized text
    """
    try:
        from nltk.corpus import stopwords
        from nltk.tokenize import sent_tokenize
        
        # Download required NLTK data if not already present
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt', quiet=True)
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords', quiet=True)
        
        # Tokenize and calculate word frequencies
        stop_words = set(stopwords.words('english'))
        sentences_list = sent_tokenize(text)
        
        # If there are fewer sentences than requested, return all
        if len(sentences_list) <= sentences:
            return text
        
        word_frequencies = {}
        for sentence in sentences_list:
            for word in nltk.word_tokenize(sentence):
                word = word.lower()
                if word not in stop_words:
                    if word not in word_frequencies:
                        word_frequencies[word] = 1
                    else:
                        word_frequencies[word] += 1
        
        # Normalize frequencies
        maximum_frequency = max(word_frequencies.values()) if word_frequencies else 1
        for word in word_frequencies:
            word_frequencies[word] = word_frequencies[word] / maximum_frequency
        
        # Score sentences
        sentence_scores = {}
        for i, sentence in enumerate(sentences_list):
            for word in nltk.word_tokenize(sentence.lower()):
                if word in word_frequencies:
                    if i not in sentence_scores:
                        sentence_scores[i] = word_frequencies[word]
                    else:
                        sentence_scores[i] += word_frequencies[word]
        
        # Get top N sentences
        summary_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)[:sentences]
        summary_sentences = [sentences_list[i] for i, _ in sorted(summary_sentences)]
        
        return ' '.join(summary_sentences)
    except Exception as e:
        logger.error(f"Summarization error: {e}")
        return text


def extract_entities(text: str) -> Dict[str, List[str]]:
    """
    Extract named entities from text.
    
    Args:
        text (str): Text to analyze
        
    Returns:
        Dict[str, List[str]]: Entities by category
    """
    try:
        import spacy
        
        # Load the spaCy model
        try:
            nlp = spacy.load("en_core_web_sm")
        except OSError:
            logger.warning("Downloading spaCy model...")
            import subprocess
            subprocess.run([sys.executable, "-m", "spacy", "download", "en_core_web_sm"], 
                           check=True)
            nlp = spacy.load("en_core_web_sm")
        
        # Process the text
        doc = nlp(text)
        
        # Extract entities
        entities = {
            "PERSON": [],
            "ORG": [],
            "GPE": [],  # Countries, cities, states
            "LOC": [],  # Non-GPE locations
            "DATE": [],
            "TIME": [],
            "MONEY": [],
            "PRODUCT": [],
            "EVENT": [],
            "WORK_OF_ART": [],
            "OTHER": []
        }
        
        for ent in doc.ents:
            if ent.label_ in entities:
                if ent.text not in entities[ent.label_]:
                    entities[ent.label_].append(ent.text)
            else:
                if ent.text not in entities["OTHER"]:
                    entities["OTHER"].append(ent.text)
        
        return entities
    except Exception as e:
        logger.error(f"Entity extraction error: {e}")
        return {"ERROR": [str(e)]}


def analyze_text(text: str) -> Dict[str, Any]:
    """
    Perform basic text analytics.
    
    Args:
        text (str): Input text
        
    Returns:
        Dict: Analysis results
    """
    try:
        # Tokenize text
        words = nltk.word_tokenize(text.lower())
        sentences = nltk.sent_tokenize(text)
        
        # Filter out punctuation
        words = [word for word in words if word.isalpha()]
        
        # Count word frequencies
        word_freq = Counter(words)
        
        # Calculate readability metrics
        avg_word_length = sum(len(word) for word in words) / len(words) if words else 0
        avg_sent_length = len(words) / len(sentences) if sentences else 0
        
        # Detect language
        language = detect_language(text)
        
        return {
            "word_count": len(words),
            "sentence_count": len(sentences),
            "unique_words": len(set(words)),
            "avg_word_length": avg_word_length,
            "avg_sentence_length": avg_sent_length,
            "language": language,
            "most_common_words": word_freq.most_common(20)
        }
    except Exception as e:
        logger.error(f"Text analysis error: {e}")
        return {"error": str(e)}


# Support for additional file types

def text_from_epub(epub_path: str) -> Optional[str]:
    """
    Extract text from EPUB ebooks.
    
    Args:
        epub_path (str): Path to the EPUB file
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        from ebooklib import epub
        import html2text
        
        book = epub.read_epub(epub_path)
        h = html2text.HTML2Text()
        h.ignore_links = False
        
        content = []
        
        # Get book metadata
        metadata = []
        if book.get_metadata('DC', 'title'):
            metadata.append(f"Title: {book.get_metadata('DC', 'title')[0][0]}")
        if book.get_metadata('DC', 'creator'):
            metadata.append(f"Author: {book.get_metadata('DC', 'creator')[0][0]}")
        if book.get_metadata('DC', 'description'):
            metadata.append(f"Description: {book.get_metadata('DC', 'description')[0][0]}")
        
        if metadata:
            content.append("\n".join(metadata))
            content.append("---")
        
        # Get book content
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                content.append(h.handle(item.get_content().decode('utf-8')))
        
        return normalize("\n".join(content))
    except ImportError:
        logger.error("ebooklib and/or html2text not installed")
        return "ebooklib and/or html2text packages are required for EPUB processing"
    except Exception as e:
        logger.error(f"Error processing EPUB: {e}")
        return None


def text_from_pptx(pptx_path: str) -> Optional[str]:
    """
    Extract text from PowerPoint presentations.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        text = ["--- PowerPoint Presentation ---"]
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {i}:"]
            
            # Get slide title if it exists
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # Extract text from all shapes
            shape_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    shape_text.append(shape.text)
            
            if shape_text:
                slide_text.append("\n".join(shape_text))
            
            text.append("\n".join(slide_text))
        
        return normalize("\n\n".join(text))
    except ImportError:
        logger.error("python-pptx not installed")
        return "python-pptx package is required for PowerPoint processing"
    except Exception as e:
        logger.error(f"Error processing PowerPoint: {e}")
        return None


def text_from_odt(odt_path: str) -> Optional[str]:
    """
    Extract text from OpenDocument Text files.
    
    Args:
        odt_path (str): Path to the ODT file
        
    Returns:
        Optional[str]: Extracted text
    """
    try:
        from odf import text, teletype
        from odf.opendocument import load
        
        textdoc = load(odt_path)
        
        # Extract metadata
        meta = []
        meta_elem = textdoc.meta
        if meta_elem:
            for prop in meta_elem.childNodes:
                if hasattr(prop, 'tagName') and hasattr(prop, 'childNodes') and prop.childNodes:
                    meta.append(f"{prop.tagName}: {teletype.extractText(prop)}")
        
        # Extract content
        allparas = textdoc.getElementsByType(text.P)
        content = "\n".join(teletype.extractText(p) for p in allparas)
        
        # Combine metadata and content
        if meta:
            final_text = "\n".join(meta) + "\n---\n" + content
        else:
            final_text = content
        
        return normalize(final_text)
    except ImportError:
        logger.error("odfpy not installed")
        return "odfpy package is required for ODT processing"
    except Exception as e:
        logger.error(f"Error processing ODT: {e}")
        return None


# Memory optimization functions

def extract_text_chunked(file_path: str, chunk_size: int = 10) -> Optional[str]:
    """
    Extract text from very large files in chunks to manage memory usage.
    
    Args:
        file_path (str): Path to the large file
        chunk_size (int): Number of pages/sections to process at once
        
    Returns:
        Optional[str]: Combined text
    """
    path = clean_path(file_path)
    if not path:
        return None
        
    mime_type = magic.from_file(path, mime=True)
    
    if mime_type == "application/pdf":
        return extract_pdf_chunked(path, chunk_size)
    elif mime_type.startswith("text/"):
        return extract_text_file_chunked(path, chunk_size)
    else:
        return extract_text(path)  # Fall back to regular extraction


def extract_text_file_chunked(file_path: str, chunk_size: int = 1024*1024) -> str:
    """
    Extract text from large text files in chunks.
    
    Args:
        file_path (str): Path to the text file
        chunk_size (int): Size of each chunk in bytes
        
    Returns:
        str: Combined text
    """
    try:
        result = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            while True:
                chunk = f.read(chunk_size)
                if not chunk:
                    break
                result.append(chunk)
        
        return normalize("".join(result))
    except Exception as e:
        logger.error(f"Error processing text file in chunks: {e}")
        return ""


# Batch processing

def batch_extract(file_paths: List[str], max_workers: Optional[int] = None) -> Dict[str, str]:
    """
    Process multiple files concurrently.
    
    Args:
        file_paths (List[str]): List of file paths to process
        max_workers (int, optional): Maximum number of worker threads
        
    Returns:
        Dict[str, str]: Dictionary mapping file paths to extracted text
    """
    results = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {executor.submit(extract_text, file_path): file_path for file_path in file_paths}
        for future in concurrent.futures.as_completed(future_to_file):
            file_path = future_to_file[future]
            try:
                results[file_path] = future.result() or f"No text extracted from {file_path}"
            except Exception as exc:
                results[file_path] = f"ERROR: {exc}"
    return results


# Enhanced CLI with more options
def main() -> None:
    """
    Enhanced CLI entry point for text extraction, metadata, and analytics.
    """
    import argparse
    import glob
    
    class TranslateAction(argparse.Action):
        def __call__(self, parser, namespace, values, option_string=None):
            # If no value is provided, set to 'list' to trigger language listing
            if values is None:
                setattr(namespace, self.dest, 'list')
            else:
                setattr(namespace, self.dest, values)
    
    parser = argparse.ArgumentParser(
        description="Extract and analyze text from any file, URL, directory or wildcard pattern"
    )
    parser.add_argument(
        "source",
        nargs="*",  # Make source completely optional
        help="Path(s) to file(s), URL, directory, or wildcard pattern"
    )
    parser.add_argument(
        "--metadata",
        action="store_true",
        help="Extract metadata instead of text"
    )
    parser.add_argument(
        "--summarize",
        action="store_true",
        help="Summarize the extracted text"
    )
    parser.add_argument(
        "--sentences",
        type=int,
        default=5,
        help="Number of sentences in summary (default: 5)"
    )
    parser.add_argument(
        "--entities",
        action="store_true",
        help="Extract named entities"
    )
    parser.add_argument(
        "--analyze",
        action="store_true",
        help="Perform text analysis"
    )
    parser.add_argument(
        "--translate",
        action=TranslateAction,
        nargs="?",  # Allow --translate with no argument to list languages
        metavar="LANG",
        help="Translate text to specified language code (e.g., 'es'), or list available languages if no code provided"
    )
    parser.add_argument(
        "--output",
        type=str,
        help="Output file path (default: stdout)"
    )
    parser.add_argument(
        "--password",
        type=str,
        help="Password for protected documents"
    )
    parser.add_argument(
        "--scrape",
        action="store_true",
        help="Scrape multiple pages from a website (for URLs only)"
    )
    parser.add_argument(
        "--max-pages",
        type=int,
        default=5,
        help="Maximum pages to scrape when using --scrape (default: 5)"
    )
    parser.add_argument(
        "--verbose", "-v",
        action="count",
        default=0,
        help="Increase verbosity (can be used multiple times)"
    )
    parser.add_argument(
        "--screenshot",
        action="store_true",
        help="Capture and extract text from screen"
    )
    parser.add_argument(
        "--chunked",
        action="store_true",
        help="Process large files in chunks to reduce memory usage"
    )
    parser.add_argument(
        "--no-js",
        action="store_true",
        help="Disable JavaScript rendering for web pages"
    )
    parser.add_argument(
        "--list-languages",
        action="store_true",
        help="List available translation languages"
    )
    
    args = parser.parse_args()
    
    # Configure logging level based on verbosity
    if args.verbose == 0:
        logger.setLevel(logging.WARNING)
    elif args.verbose == 1:
        logger.setLevel(logging.INFO)
    else:
        logger.setLevel(logging.DEBUG)
    
    # Handle translation language listing (in two ways for better usability)
    if args.translate == "list" or args.list_languages:
        try:
            logger.info("Fetching available languages...")
            languages = list_available_languages()
            
            print("\n[Available Translation Languages]")
            # Format the output as a table with columns
            max_code_len = max(len(code) for code in languages.keys())
            max_name_len = max(len(name) for name in languages.values())
            format_str = f"  {{:<{max_code_len}}}  {{:<{max_name_len}}}"
            
            print(format_str.format("Code", "Language"))
            print(format_str.format("-" * max_code_len, "-" * max_name_len))
            
            for code, name in sorted(languages.items(), key=lambda x: x[1]):  # Sort by language name
                print(format_str.format(code, name))
            
            sys.exit(0)
        except Exception as e:
            logger.error(f"Error listing languages: {e}")
            sys.exit(1)
    
    # Ensure source is provided for all other operations
    if not args.source and not args.screenshot:
        parser.print_help()
        sys.exit()
    
    # Function to handle output
    def output_result(result):
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                if isinstance(result, dict):
                    json.dump(result, f, indent=2, ensure_ascii=False)
                else:
                    f.write(str(result))
            logger.info(f"Output written to {args.output}")
        else:
            if isinstance(result, dict):
                print(json.dumps(result, indent=2, ensure_ascii=False))
            else:
                print(result)
    
    # Handle screenshot mode
    if args.screenshot:
        logger.info("Capturing screenshot...")
        text = text_from_screenshot()
        if not text:
            logger.error("No text extracted from screenshot.")
            sys.exit(1)
        
        # Process the text according to arguments
        result = process_extracted_text(text, args)
        output_result(result)
        sys.exit(0)
    
    # Process sources
    sources = args.source
    
    # Expand any directory sources
    files = []
    for source in sources:
        if os.path.isdir(source):
            # It's a directory, add all files in it
            logger.info(f"Processing directory: {source}")
            dir_files = [os.path.join(source, f) for f in os.listdir(source) 
                        if os.path.isfile(os.path.join(source, f))]
            files.extend(dir_files)
        elif any(c in source for c in ["*", "?", "[", "]"]) and not os.path.exists(source):
            # It's a wildcard pattern that wasn't expanded by the shell
            logger.info(f"Processing pattern: {source}")
            expanded = glob.glob(source, recursive=True)
            if expanded:
                files.extend(expanded)
            else:
                logger.warning(f"No files found matching pattern: {source}")
        elif os.path.isfile(source):
            # It's a regular file
            files.append(source)
        elif is_url(source):
            # It's a URL - we'll handle it separately
            if len(sources) == 1:
                # Only process URL if it's the only source
                if args.scrape:
                    logger.info(f"Scraping website: {source} (max {args.max_pages} pages)")
                    results = scrape_website(source, max_pages=args.max_pages)
                    combined_text = "\n\n".join([f"=== {url} ===\n{text}" for url, text in results.items()])
                    result = process_extracted_text(combined_text, args)
                    output_result(result)
                else:
                    logger.info(f"Extracting text from URL: {source}")
                    text = text_from_url(source, render_js=not args.no_js)
                    if not text:
                        logger.error("No text extracted.")
                        sys.exit(1)
                    result = process_extracted_text(text, args)
                    output_result(result)
                return
            else:
                logger.warning(f"Skipping URL {source} when processing multiple sources")
        else:
            logger.warning(f"Source not found: {source}")
    
    if not files:
        logger.error("No valid files found to process.")
        sys.exit(1)
    
    logger.info(f"Processing {len(files)} files")
    
    # Process metadata if requested
    if args.metadata:
        results = {}
        with concurrent.futures.ThreadPoolExecutor() as executor:
            future_to_file = {executor.submit(extract_metadata, file_path): file_path for file_path in files}
            for future in concurrent.futures.as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    results[os.path.basename(file_path)] = future.result()
                except Exception as exc:
                    results[os.path.basename(file_path)] = {"error": str(exc)}
        output_result(results)
        return
    
    # If only one file and not metadata, process it directly
    if len(files) == 1:
        file_path = files[0]
        try:
            if args.chunked:
                logger.info(f"Extracting text from {file_path} in chunks")
                text = extract_text_chunked(file_path)
            elif args.password:
                logger.info(f"Extracting text from password-protected {file_path}")
                text = extract_text_with_password(file_path, args.password)
            else:
                logger.info(f"Extracting text from {file_path}")
                text = extract_text(file_path)
                
            if not text:
                logger.error("No text extracted.")
                sys.exit(1)
                
            result = process_extracted_text(text, args)
            output_result(result)
            return
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            sys.exit(1)
    
    # Process multiple files
    results = batch_extract(files)
    
    # Process each extracted text according to arguments
    processed_results = {}
    for file_path, text in results.items():
        if isinstance(text, str) and not text.startswith("ERROR:"):
            processed_results[os.path.basename(file_path)] = process_extracted_text(text, args)
        else:
            processed_results[os.path.basename(file_path)] = text
    
    output_result(processed_results)


def process_extracted_text(text: str, args) -> Union[str, Dict[str, Any]]:
    """
    Process extracted text according to command-line arguments.
    
    Args:
        text (str): Extracted text
        args: Command-line arguments
        
    Returns:
        Union[str, Dict[str, Any]]: Processed text or analysis results
    """
    # Process text according to arguments
    if args.summarize:
        logger.info(f"Summarizing text (sentences: {args.sentences})")
        text = summarize_text(text, args.sentences)
    
    if args.entities:
        logger.info("Extracting named entities")
        entities = extract_entities(text)
        return entities
    
    if args.analyze:
        logger.info("Analyzing text")
        analysis = analyze_text(text)
        return analysis
    
    if args.translate and args.translate != "list":
        logger.info(f"Translating text to {args.translate}")
        translated = translate_text(text, args.translate)
        if translated:
            text = translated
        else:
            logger.error("Translation failed")
    
    return text


# Simple REST API server
def start_api_server(host: str = 'localhost', port: int = 8000) -> None:
    """
    Start a simple API server to expose the text extraction functionality.
    
    Args:
        host (str): Host to bind to
        port (int): Port to listen on
    """
    try:
        from flask import Flask, request, jsonify
        
        app = Flask(__name__)
        
        @app.route('/extract', methods=['POST'])
        def extract_endpoint():
            if 'file' in request.files:
                file = request.files['file']
                temp_path = os.path.join(tempfile.gettempdir(), f"upload_{uuid4().hex}")
                file.save(temp_path)
                
                try:
                    params = request.form.to_dict()
                    
                    # Extract text
                    if params.get('password'):
                        text = extract_text_with_password(temp_path, params['password'])
                    else:
                        text = extract_text(temp_path)
                        
                    if not text:
                        return jsonify({"error": "No text extracted"}), 400
                    
                    # Process text according to parameters
                    result = {"text": text}
                    
                    if params.get('summarize') == 'true':
                        sentences = int(params.get('sentences', 5))
                        result["summary"] = summarize_text(text, sentences)
                        
                    if params.get('entities') == 'true':
                        result["entities"] = extract_entities(text)
                        
                    if params.get('analyze') == 'true':
                        result["analysis"] = analyze_text(text)
                        
                    if params.get('translate'):
                        translated = translate_text(text, params['translate'])
                        if translated:
                            result["translated"] = translated
                    
                    return jsonify(result)
                finally:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
            
            elif 'url' in request.json:
                url = request.json['url']
                params = request.json
                
                if params.get('scrape') == True:
                    max_pages = int(params.get('max_pages', 5))
                    results = scrape_website(url, max_pages=max_pages)
                    combined_text = "\n\n".join([f"=== {url} ===\n{text}" for url, text in results.items()])
                    result = {"text": combined_text, "pages_scraped": len(results)}
                else:
                    text = text_from_url(url)
                    if not text:
                        return jsonify({"error": "No text extracted from URL"}), 400
                    result = {"text": text}
                
                # Process text according to parameters
                if params.get('summarize') == True:
                    sentences = int(params.get('sentences', 5))
                    result["summary"] = summarize_text(result["text"], sentences)
                    
                if params.get('entities') == True:
                    result["entities"] = extract_entities(result["text"])
                    
                if params.get('analyze') == True:
                    result["analysis"] = analyze_text(result["text"])
                    
                if params.get('translate'):
                    translated = translate_text(result["text"], params['translate'])
                    if translated:
                        result["translated"] = translated
                
                return jsonify(result)
            
            return jsonify({"error": "No file or URL provided"}), 400
        
        @app.route('/metadata', methods=['POST'])
        def metadata_endpoint():
            if 'file' in request.files:
                file = request.files['file']
                temp_path = os.path.join(tempfile.gettempdir(), f"upload_{uuid4().hex}")
                file.save(temp_path)
                
                try:
                    metadata = extract_metadata(temp_path)
                    return jsonify(metadata)
                finally:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
            
            elif 'url' in request.json:
                url = request.json['url']
                try:
                    # For URLs, we can provide limited metadata
                    import requests
                    resp = requests.head(url)
                    metadata = {
                        "url": url,
                        "status_code": resp.status_code,
                        "content_type": resp.headers.get('Content-Type'),
                        "content_length": resp.headers.get('Content-Length'),
                        "server": resp.headers.get('Server'),
                        "last_modified": resp.headers.get('Last-Modified'),
                    }
                    return jsonify(metadata)
                except Exception as e:
                    return jsonify({"error": str(e)}), 400
        
            return jsonify({"error": "No file or URL provided"}), 400
        
        @app.route('/screenshot', methods=['GET'])
        def screenshot_endpoint():
            try:
                text = text_from_screenshot()
                if not text:
                    return jsonify({"error": "No text extracted from screenshot"}), 400
                
                result = {"text": text}
                params = request.args.to_dict()
                
                # Process text according to parameters
                if params.get('summarize') == 'true':
                    sentences = int(params.get('sentences', 5))
                    result["summary"] = summarize_text(text, sentences)
                    
                if params.get('entities') == 'true':
                    result["entities"] = extract_entities(text)
                    
                if params.get('analyze') == 'true':
                    result["analysis"] = analyze_text(text)
                
                return jsonify(result)
            except Exception as e:
                return jsonify({"error": str(e)}), 500
        
        logger.info(f"Starting API server on {host}:{port}")
        app.run(host=host, port=port)
    except ImportError:
        logger.error("Flask not installed. Install with: pip install flask")
        

# Execute as script
if __name__ == "__main__":
    main()    
